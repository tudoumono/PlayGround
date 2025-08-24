# ==============================================================================
# Universal Python Design Ver. 2.0
#
# このスクリプトは、指定されたslide_dataに基づいて、
# Google風デザインに準拠したPowerPointプレゼンテーションを自動生成します。
# 必要なライブラリ: pip install python-pptx Pillow requests
# ==============================================================================
import io
import re
import requests
from pathlib import Path
from datetime import date
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

# ==============================================================================
# 1. マスターデザイン設定 (このセクションは変更しないこと)
# ==============================================================================
SETTINGS = {
    'OUTPUT_FILENAME': 'presentation_extended.pptx'
}

CONFIG = {
    'BASE_PX': {'W': 960, 'H': 540},
    'BASE_EMU': {'W': Inches(10).emu, 'H': Inches(5.625).emu},
    'POS_PX': {
        'titleSlide': {
            'logo':      {'left': 55,  'top': 105, 'width': 135},
            'title':     {'left': 50,  'top': 230, 'width': 800, 'height': 90},
            'date':      {'left': 50,  'top': 340, 'width': 250, 'height': 40},
        },
        'contentSlide': {
            'headerLogo':    {'right': 20, 'top': 20, 'width': 75},
            'title':         {'left': 25, 'top': 60,  'width': 830, 'height': 65},
            'titleUnderline':{'left': 25, 'top': 128, 'width': 260, 'height': 4},
            'subhead':       {'left': 25, 'top': 140, 'width': 830, 'height': 30},
            'body':          {'left': 25, 'top': 172, 'width': 910, 'height': 303},
            'gridArea':      {'left': 25, 'top': 172, 'width': 910, 'height': 303},
            'compareLeft':   {'left': 25, 'top': 172, 'width': 430, 'height': 303},
            'compareRight':  {'left': 505, 'top': 172, 'width': 430, 'height': 303},
        },
        'sectionSlide': {
            'title':    {'left': 55, 'top': 230, 'width': 840, 'height': 80},
            'ghostNum': {'left': 35, 'top': 120, 'width': 300, 'height': 200},
        },
        'footer': {
            'leftText':  {'left': 15, 'top': 505, 'width': 250, 'height': 20},
            'rightPage': {'right': 15, 'top': 505, 'width': 50,  'height': 20},
        },
        'bottomBar': {'left': 0, 'top': 534, 'width': 960, 'height': 6}
    },
    'FONTS': {
        'family': 'Arial',
        'sizes': {
            'title': 45, 'date': 16, 'sectionTitle': 38, 'contentTitle': 28,
            'subhead': 18, 'body': 14, 'footer': 9, 'cardTitle': 16,
            'cardDesc': 12, 'ghostNum': 180, 'processStep': 14, 'timelineLabel': 10,
            'laneTitle': 13,
        }
    },
    'COLORS': {
        'primary_blue': '4285F4', 'google_red': 'EA4335', 'google_yellow': 'FBBC04',
        'google_green': '34A853', 'text_primary': '333333', 'text_white': 'FFFFFF',
        'background_white': 'FFFFFF', 'background_gray': 'F8F9FA', 'card_bg': 'FFFFFF',
        'card_border': 'DADCE0', 'ghost_gray': 'EFEFED', 'faint_gray': 'E8EAED',
        'neutral_gray': '9E9E9E', 'lane_title_bg': 'F5F5F3', 'lane_border': 'DADCE0',
    },
    'LOGOS': {
        'header': 'https://upload.wikimedia.org/wikipedia/commons/thumb/2/2f/Google_2015_logo.svg/1024px-Google_2015_logo.svg.png',
        'closing': 'https://upload.wikimedia.org/wikipedia/commons/thumb/2/2f/Google_2015_logo.svg/1024px-Google_2015_logo.svg.png'
    },
    'FOOTER_TEXT': f"© {date.today().year} Your Organization"
}

# ==============================================================================
# 2. スライドデータ (このセクションはAIによって生成・置換される)
# ==============================================================================
slide_data = [
    {'type': 'title', 'title': 'CDSによる統合文書管理・活用', 'date': '2025.08.21', 'notes': '本日は、効率的な文書管理・活用を実現する「コーポレート・ドキュメント・システム」、通称CDSについてご説明します。'},
    {'type': 'content', 'title': 'アジェンダ', 'points': ['1. CDS（コーポレート・ドキュメント・システム）とは', '2. CDSを構成する3つの要素', '3. まとめ'], 'notes': '本日はこちらの流れでご説明いたします。まずCDSの全体像を掴んでいただき、次にその中核となる3つの要素について詳しく見ていきます。'},
    {'type': 'section', 'title': '1. CDS（コーポレート・ドキュメント・システム）とは', 'notes': '最初のセクションでは、我々が提唱する統合的な文書管理の概念「CDS」についてご説明します。'},
    {'type': 'content', 'title': '統合的な文書管理・活用システム「CDS」', 'subhead': '文書管理の課題を「標準化」「電子化」「共有化」の3つの観点から統合的に解決する仕組み', 'points': ['文書に[[統一性]]をもたせる「標準化」', '作成・伝達を[[容易]]にする「電子化」', '誰でも[[利用可能]]にする「共有化」'], 'notes': 'CDSは、単なるツールの導入ではありません。ここに挙げた「標準化」「電子化」「共有化」という3つの柱を統合的に進めることで、組織全体の知的生産性を向上させることを目的としたシステムです。'},
    {'type': 'section', 'title': '2. CDSを構成する3つの要素', 'notes': '次のセクションでは、CDSを成り立たせている3つの要素、標準化、電子化、共有化について、それぞれ掘り下げて見ていきましょう。'},
    {'type': 'cards', 'title': 'CDSを支える3つの柱', 'subhead': 'それぞれの観点が持つ役割と目的を理解する', 'columns': 3, 'items': [{'title': '標準化 (Standardization)', 'desc': '文書管理、デザイン、入力のルールを定め、文書の**統一性**を持たせる'}, {'title': '電子化 (Digitization)', 'desc': '文書メディアを紙から電子に移行し、作成・生産・伝達を**容易**にする'}, {'title': '共有化 (Sharing)', 'desc': 'DBやネットワークを活用し、必要な文書を**誰もが利用できる**状態にする'}], 'notes': 'こちらが3つの柱の概要です。標準化で「質」を担保し、電子化で「効率」を高め、共有化で「価値」を最大化する、という関係性になります。'},
    {'type': 'table', 'title': '3要素の具体的な取り組み', 'subhead': '各要素の目的、具体策、そして成功に必要なものを整理', 'headers': ['観点', '目的', '具体策', '成功の鍵'], 'rows': [['**標準化**', '文書の統一性', '文書管理/デザイン/入力標準の策定', 'マニュアル整備、教育'], ['**電子化**', '作成・生産・伝達の容易化', 'ペーパーレス化の推進', '十分なスペース、適切なツール'], ['**共有化**', '全社的な情報活用', 'DB/ライブラリ/ネットワークの整備', '管理・運営ルールの策定']], 'notes': 'この表は、3つの要素を具体的に進める上でのアクションアイテムをまとめたものです。例えば、標準化においてはルールを作るだけでなく、それを浸透させるためのマニュアルや教育が不可欠です。'},
    {'type': 'closing', 'notes': '以上、CDSの概念とその構成要素についてご説明しました。標準化、電子化、共有化を三位一体で進めることで、組織の文書資産を最大限に活用することが可能になります。ご清聴ありがとうございました。'}
]


# ==============================================================================
# 3. コアロジック (このセクションは変更しないこと)
# ==============================================================================
class LayoutManager:
    def __init__(self, page_w_emu, page_h_emu):
        self.page_w_emu, self.page_h_emu = page_w_emu, page_h_emu
        self.base_w_px, self.base_h_px = CONFIG['BASE_PX']['W'], CONFIG['BASE_PX']['H']
        self.scale_x, self.scale_y = self.page_w_emu / self.base_w_px, self.page_h_emu / self.base_h_px

    def px_to_emu(self, px, axis='x'):
        return int(px * (self.scale_x if axis == 'x' else self.scale_y))

    def get_rect(self, spec_path):
        keys = spec_path.split('.')
        pos_px = CONFIG['POS_PX']
        for key in keys: pos_px = pos_px[key]
        left_px = pos_px.get('left')
        if pos_px.get('right') is not None and left_px is None:
            left_px = self.base_w_px - pos_px['right'] - pos_px['width']
        return {
            'left': self.px_to_emu(left_px, 'x') if left_px is not None else None,
            'top': self.px_to_emu(pos_px['top'], 'y') if pos_px.get('top') is not None else None,
            'width': self.px_to_emu(pos_px['width'], 'x') if pos_px.get('width') is not None else None,
            'height': self.px_to_emu(pos_px['height'], 'y') if pos_px.get('height') is not None else None,
        }

class SlideGenerator:
    def __init__(self):
        self._section_counter = 0
        self.slide_generators = {
            'title': self._create_title_slide, 'section': self._create_section_slide,
            'content': self._create_content_slide, 'cards': self._create_cards_slide,
            'table': self._create_table_slide, 'closing': self._create_closing_slide,
            'compare': self._create_compare_slide, 'process': self._create_process_slide,
            'timeline': self._create_timeline_slide, 'diagram': self._create_diagram_slide,
            'progress': self._create_progress_slide,
        }

    def _set_font_style(self, run, style_opts):
        font = run.font
        font.name = style_opts.get('family', CONFIG['FONTS']['family'])
        font.size = Pt(style_opts.get('size', CONFIG['FONTS']['sizes']['body']))
        font.bold = style_opts.get('bold', False)
        color_str = style_opts.get('color', CONFIG['COLORS']['text_primary'])
        font.color.rgb = RGBColor.from_string(color_str)

    def _parse_inline_styles(self, text):
        pattern = r'(\*\*|\[\[)(.*?)\1'
        parts, last_end = [], 0
        for match in re.finditer(pattern, text):
            start, end = match.span()
            tag, content = match.groups()
            if start > last_end: parts.append({'text': text[last_end:start], 'style': {}})
            style = {'bold': True, 'color': CONFIG['COLORS']['primary_blue']} if tag == '[[' else {'bold': True}
            parts.append({'text': content, 'style': style})
            last_end = end
        if last_end < len(text): parts.append({'text': text[last_end:], 'style': {}})
        return parts

    def _set_styled_text(self, text_frame, text, base_style_opts):
        text_frame.clear()
        p = text_frame.paragraphs[0]
        p.font.name = base_style_opts.get('family', CONFIG['FONTS']['family'])
        p.font.size = Pt(base_style_opts.get('size', CONFIG['FONTS']['sizes']['body']))
        p.alignment = base_style_opts.get('align', PP_ALIGN.LEFT)
        for i, line in enumerate(text.split('\n')):
            if i > 0: p = text_frame.add_paragraph()
            parts = self._parse_inline_styles(line) or [{'text': line, 'style': {}}]
            for part in parts:
                run = p.add_run()
                run.text = part['text']
                self._set_font_style(run, {**base_style_opts, **part['style']})

    def _set_bullets_with_inline_styles(self, text_frame, points, base_style_opts):
        text_frame.clear()
        for i, point_text in enumerate(points):
            p = text_frame.paragraphs[i] if i == 0 else text_frame.add_paragraph()
            p.text = "• "
            p.font.size = Pt(base_style_opts.get('size', CONFIG['FONTS']['sizes']['body']))
            p.font.name = base_style_opts.get('family', CONFIG['FONTS']['family'])
            p.level = 0
            
            parts = self._parse_inline_styles(point_text)
            for part in parts:
                run = p.add_run()
                run.text = part['text']
                self._set_font_style(run, {**base_style_opts, **part['style']})

    def _draw_bottom_bar(self, slide, layout):
        rect = layout.get_rect('bottomBar')
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, rect['left'], rect['top'], rect['width'], rect['height'])
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor.from_string(CONFIG['COLORS']['primary_blue'])
        shape.line.fill.background()

    def _add_google_footer(self, slide, layout, page_num):
        rect = layout.get_rect('footer.leftText')
        textbox = slide.shapes.add_textbox(rect['left'], rect['top'], rect['width'], rect['height'])
        self._set_styled_text(textbox.text_frame, CONFIG['FOOTER_TEXT'], {'size': CONFIG['FONTS']['sizes']['footer']})
        if page_num > 0:
            rect = layout.get_rect('footer.rightPage')
            textbox = slide.shapes.add_textbox(rect['left'], rect['top'], rect['width'], rect['height'])
            self._set_styled_text(textbox.text_frame, str(page_num), {
                'size': CONFIG['FONTS']['sizes']['footer'], 'color': CONFIG['COLORS']['primary_blue'], 'align': PP_ALIGN.RIGHT
            })

    def _draw_standard_title_header(self, slide, layout, key, title):
        rect = layout.get_rect(f'{key}.headerLogo')
        try:
            response = requests.get(CONFIG['LOGOS']['header'])
            slide.shapes.add_picture(io.BytesIO(response.content), rect['left'], rect['top'], width=rect['width'])
        except Exception: pass
        rect = layout.get_rect(f'{key}.title')
        textbox = slide.shapes.add_textbox(rect['left'], rect['top'], rect['width'], rect['height'])
        self._set_styled_text(textbox.text_frame, title, {'size': CONFIG['FONTS']['sizes']['contentTitle'], 'bold': True})
        rect = layout.get_rect(f'{key}.titleUnderline')
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, rect['left'], rect['top'], rect['width'], rect['height'])
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor.from_string(CONFIG['COLORS']['primary_blue'])
        shape.line.fill.background()

    def _draw_subhead_if_any(self, slide, layout, key, subhead):
        if not subhead: return 0
        rect = layout.get_rect(f'{key}.subhead')
        textbox = slide.shapes.add_textbox(rect['left'], rect['top'], rect['width'], rect['height'])
        self._set_styled_text(textbox.text_frame, subhead, {'size': CONFIG['FONTS']['sizes']['subhead']})
        return layout.px_to_emu(36)

    def _create_title_slide(self, slide, data, layout, page_num):
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor.from_string(CONFIG['COLORS']['background_white'])
        rect = layout.get_rect('titleSlide.logo')
        try:
            response = requests.get(CONFIG['LOGOS']['header'])
            slide.shapes.add_picture(io.BytesIO(response.content), rect['left'], rect['top'], width=rect['width'])
        except Exception: pass
        rect = layout.get_rect('titleSlide.title')
        textbox = slide.shapes.add_textbox(rect['left'], rect['top'], rect['width'], rect['height'])
        self._set_styled_text(textbox.text_frame, data.get('title', ''), {'size': CONFIG['FONTS']['sizes']['title'], 'bold': True})
        rect = layout.get_rect('titleSlide.date')
        textbox = slide.shapes.add_textbox(rect['left'], rect['top'], rect['width'], rect['height'])
        self._set_styled_text(textbox.text_frame, data.get('date', ''), {'size': CONFIG['FONTS']['sizes']['date']})
        self._draw_bottom_bar(slide, layout)

    def _create_section_slide(self, slide, data, layout, page_num):
        self._section_counter += 1
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor.from_string(CONFIG['COLORS']['background_gray'])
        num_str = str(data.get('sectionNo', self._section_counter)).zfill(2)
        rect = layout.get_rect('sectionSlide.ghostNum')
        textbox = slide.shapes.add_textbox(rect['left'], rect['top'], rect['width'], rect['height'])
        self._set_styled_text(textbox.text_frame, num_str, {
            'size': CONFIG['FONTS']['sizes']['ghostNum'], 'bold': True, 'color': CONFIG['COLORS']['ghost_gray'], 'align': PP_ALIGN.CENTER
        })
        rect = layout.get_rect('sectionSlide.title')
        textbox = slide.shapes.add_textbox(rect['left'], rect['top'], rect['width'], rect['height'])
        self._set_styled_text(textbox.text_frame, data.get('title', ''), {
            'size': CONFIG['FONTS']['sizes']['sectionTitle'], 'bold': True, 'align': PP_ALIGN.CENTER
        })
        self._add_google_footer(slide, layout, page_num)

    def _create_content_slide(self, slide, data, layout, page_num):
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor.from_string(CONFIG['COLORS']['background_white'])
        self._draw_standard_title_header(slide, layout, 'contentSlide', data.get('title'))
        dy = self._draw_subhead_if_any(slide, layout, 'contentSlide', data.get('subhead'))
        body_rect = layout.get_rect('contentSlide.body')
        body_rect['top'] += dy
        textbox = slide.shapes.add_textbox(body_rect['left'], body_rect['top'], body_rect['width'], body_rect['height'])
        self._set_bullets_with_inline_styles(textbox.text_frame, data.get('points', []), {})
        self._draw_bottom_bar(slide, layout)
        self._add_google_footer(slide, layout, page_num)

    def _create_cards_slide(self, slide, data, layout, page_num):
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor.from_string(CONFIG['COLORS']['background_white'])
        self._draw_standard_title_header(slide, layout, 'contentSlide', data.get('title'))
        dy = self._draw_subhead_if_any(slide, layout, 'contentSlide', data.get('subhead'))
        area = layout.get_rect('contentSlide.gridArea')
        area['top'] += dy
        items = data.get('items', [])
        cols = data.get('columns', 3)
        rows = -(-len(items) // cols)
        gap = layout.px_to_emu(16)
        card_w = (area['width'] - gap * (cols - 1)) // cols
        card_h = (area['height'] - gap * (rows - 1)) // rows
        for i, item in enumerate(items):
            r, c = divmod(i, cols)
            left, top = area['left'] + c * (card_w + gap), area['top'] + r * (card_h + gap)
            shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, card_w, card_h)
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor.from_string(CONFIG['COLORS']['card_bg'])
            shape.line.color.rgb = RGBColor.from_string(CONFIG['COLORS']['card_border'])
            shape.line.width = Pt(1).emu
            tf = shape.text_frame
            tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = layout.px_to_emu(10)
            tf.word_wrap = True
            tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
            if isinstance(item, dict):
                self._set_styled_text(tf, f"**{item.get('title', '')}**\n{item.get('desc', '')}", {'size': CONFIG['FONTS']['sizes']['cardDesc']})
        self._draw_bottom_bar(slide, layout)
        self._add_google_footer(slide, layout, page_num)

    def _create_table_slide(self, slide, data, layout, page_num):
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor.from_string(CONFIG['COLORS']['background_white'])
        self._draw_standard_title_header(slide, layout, 'contentSlide', data.get('title'))
        dy = self._draw_subhead_if_any(slide, layout, 'contentSlide', data.get('subhead'))
        area = layout.get_rect('contentSlide.body')
        area['top'] += dy
        headers = data.get('headers', [])
        rows_data = data.get('rows', [])
        if not headers or not rows_data: return
        shape = slide.shapes.add_table(len(rows_data) + 1, len(headers), area['left'], area['top'], area['width'], area['height'])
        table = shape.table
        for c, header_text in enumerate(headers):
            cell = table.cell(0, c)
            cell.text = header_text
            cell.text_frame.paragraphs[0].font.bold = True
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        for r, row_data in enumerate(rows_data):
            for c, cell_text in enumerate(row_data):
                cell = table.cell(r + 1, c)
                cell.text = cell_text
                cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        self._draw_bottom_bar(slide, layout)
        self._add_google_footer(slide, layout, page_num)

    def _create_compare_slide(self, slide, data, layout, page_num):
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor.from_string(CONFIG['COLORS']['background_white'])
        self._draw_standard_title_header(slide, layout, 'contentSlide', data.get('title'))
        dy = self._draw_subhead_if_any(slide, layout, 'contentSlide', data.get('subhead'))
        
        def draw_compare_box(rect, title, items):
            box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, rect['left'], rect['top'], rect['width'], rect['height'])
            box.fill.solid()
            box.fill.fore_color.rgb = RGBColor.from_string(CONFIG['COLORS']['lane_title_bg'])
            box.line.color.rgb = RGBColor.from_string(CONFIG['COLORS']['lane_border'])
            box.line.width = Pt(1).emu
            
            title_h = layout.px_to_emu(40)
            title_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, rect['left'], rect['top'], rect['width'], title_h)
            title_bar.fill.solid()
            title_bar.fill.fore_color.rgb = RGBColor.from_string(CONFIG['COLORS']['primary_blue'])
            title_bar.line.fill.background()
            title_bar.text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
            self._set_styled_text(title_bar.text_frame, title, {'size': CONFIG['FONTS']['sizes']['laneTitle'], 'bold': True, 'color': CONFIG['COLORS']['text_white'], 'align': PP_ALIGN.CENTER})

            pad = layout.px_to_emu(12)
            text_box = slide.shapes.add_textbox(rect['left'] + pad, rect['top'] + title_h + pad, rect['width'] - 2*pad, rect['height'] - title_h - 2*pad)
            self._set_bullets_with_inline_styles(text_box.text_frame, items, {})

        left_rect = layout.get_rect('contentSlide.compareLeft')
        left_rect['top'] += dy
        draw_compare_box(left_rect, data.get('leftTitle', ''), data.get('leftItems', []))
        
        right_rect = layout.get_rect('contentSlide.compareRight')
        right_rect['top'] += dy
        draw_compare_box(right_rect, data.get('rightTitle', ''), data.get('rightItems', []))

        self._draw_bottom_bar(slide, layout)
        self._add_google_footer(slide, layout, page_num)

    def _create_process_slide(self, slide, data, layout, page_num):
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor.from_string(CONFIG['COLORS']['background_white'])
        self._draw_standard_title_header(slide, layout, 'contentSlide', data.get('title'))
        dy = self._draw_subhead_if_any(slide, layout, 'contentSlide', data.get('subhead'))
        
        area = layout.get_rect('contentSlide.body')
        area['top'] += dy
        steps = data.get('steps', [])
        n = len(steps)
        if n == 0: return

        gap_y = (area['height'] - layout.px_to_emu(40)) / max(1, n - 1) if n > 1 else 0
        cx = area['left'] + layout.px_to_emu(44)
        top0 = area['top'] + layout.px_to_emu(20)

        if n > 1:
            line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, cx - layout.px_to_emu(1), top0, layout.px_to_emu(2), gap_y * (n - 1))
            line.fill.solid()
            line.fill.fore_color.rgb = RGBColor.from_string(CONFIG['COLORS']['faint_gray'])
            line.line.fill.background()

        for i, step_text in enumerate(steps):
            cy = top0 + gap_y * i
            sz = layout.px_to_emu(28)
            num_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, cx - sz/2, cy - sz/2, sz, sz)
            num_box.fill.solid()
            num_box.fill.fore_color.rgb = RGBColor.from_string(CONFIG['COLORS']['background_white'])
            num_box.line.color.rgb = RGBColor.from_string(CONFIG['COLORS']['primary_blue'])
            num_box.line.width = Pt(1).emu
            num_box.text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
            self._set_styled_text(num_box.text_frame, str(i + 1), {'size': 12, 'bold': True, 'color': CONFIG['COLORS']['primary_blue'], 'align': PP_ALIGN.CENTER})

            txt_box = slide.shapes.add_textbox(cx + layout.px_to_emu(28), cy - layout.px_to_emu(16), area['width'] - layout.px_to_emu(100), layout.px_to_emu(32))
            txt_box.text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
            self._set_styled_text(txt_box.text_frame, step_text, {'size': CONFIG['FONTS']['sizes']['processStep']})
            
        self._draw_bottom_bar(slide, layout)
        self._add_google_footer(slide, layout, page_num)

    def _create_timeline_slide(self, slide, data, layout, page_num):
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor.from_string(CONFIG['COLORS']['background_white'])
        self._draw_standard_title_header(slide, layout, 'contentSlide', data.get('title'))
        dy = self._draw_subhead_if_any(slide, layout, 'contentSlide', data.get('subhead'))
        
        area = layout.get_rect('contentSlide.body')
        area['top'] += dy
        milestones = data.get('milestones', [])
        n = len(milestones)
        if n == 0: return

        inner_margin = layout.px_to_emu(60)
        base_y = area['top'] + area['height'] * 0.5
        left_x, right_x = area['left'] + inner_margin, area['left'] + area['width'] - inner_margin

        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left_x, base_y - layout.px_to_emu(1), right_x - left_x, layout.px_to_emu(2))
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor.from_string(CONFIG['COLORS']['faint_gray'])
        line.line.fill.background()

        dot_r = layout.px_to_emu(12)
        gap_x = (right_x - left_x) / max(1, n - 1) if n > 1 else 0

        for i, m in enumerate(milestones):
            x = left_x + gap_x * i
            dot = slide.shapes.add_shape(MSO_SHAPE.ELLIPSE, x - dot_r/2, base_y - dot_r/2, dot_r, dot_r)
            state = m.get('state', 'todo').lower()
            if state == 'done':
                dot.fill.solid()
                dot.fill.fore_color.rgb = RGBColor.from_string(CONFIG['COLORS']['google_green'])
                dot.line.fill.background()
            elif state == 'next':
                dot.fill.solid()
                dot.fill.fore_color.rgb = RGBColor.from_string(CONFIG['COLORS']['background_white'])
                dot.line.color.rgb = RGBColor.from_string(CONFIG['COLORS']['google_yellow'])
                dot.line.width = Pt(2).emu
            else:
                dot.fill.solid()
                dot.fill.fore_color.rgb = RGBColor.from_string(CONFIG['COLORS']['background_white'])
                dot.line.color.rgb = RGBColor.from_string(CONFIG['COLORS']['neutral_gray'])
                dot.line.width = Pt(1).emu
            
            label_box = slide.shapes.add_textbox(x - layout.px_to_emu(50), base_y - layout.px_to_emu(50), layout.px_to_emu(100), layout.px_to_emu(30))
            self._set_styled_text(label_box.text_frame, m.get('label', ''), {'size': CONFIG['FONTS']['sizes']['timelineLabel'], 'bold': True, 'align': PP_ALIGN.CENTER})
            
            date_box = slide.shapes.add_textbox(x - layout.px_to_emu(50), base_y + layout.px_to_emu(15), layout.px_to_emu(100), layout.px_to_emu(20))
            self._set_styled_text(date_box.text_frame, m.get('date', ''), {'size': CONFIG['FONTS']['sizes']['timelineLabel'], 'color': CONFIG['COLORS']['neutral_gray'], 'align': PP_ALIGN.CENTER})

        self._draw_bottom_bar(slide, layout)
        self._add_google_footer(slide, layout, page_num)

    def _create_diagram_slide(self, slide, data, layout, page_num):
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor.from_string(CONFIG['COLORS']['background_white'])
        self._draw_standard_title_header(slide, layout, 'contentSlide', data.get('title'))
        dy = self._draw_subhead_if_any(slide, layout, 'contentSlide', data.get('subhead'))
        
        area = layout.get_rect('contentSlide.body')
        area['top'] += dy
        lanes = data.get('lanes', [])
        n = len(lanes)
        if n == 0: return

        lane_gap = layout.px_to_emu(24)
        lane_pad = layout.px_to_emu(10)
        lane_title_h = layout.px_to_emu(30)
        card_gap = layout.px_to_emu(12)
        card_min_h = layout.px_to_emu(48)
        arrow_h = layout.px_to_emu(10)
        arrow_gap = layout.px_to_emu(8)

        lane_w = (area['width'] - lane_gap * (n - 1)) / n
        card_boxes = [[] for _ in range(n)]

        for j, lane in enumerate(lanes):
            left = area['left'] + j * (lane_w + lane_gap)
            
            lt = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, area['top'], lane_w, lane_title_h)
            lt.fill.solid()
            lt.fill.fore_color.rgb = RGBColor.from_string(CONFIG['COLORS']['lane_title_bg'])
            lt.line.color.rgb = RGBColor.from_string(CONFIG['COLORS']['lane_border'])
            lt.text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
            self._set_styled_text(lt.text_frame, lane.get('title', ''), {'size': CONFIG['FONTS']['sizes']['laneTitle'], 'bold': True, 'align': PP_ALIGN.CENTER})
            
            items = lane.get('items', [])
            rows = len(items)
            avail_h = area['height'] - lane_title_h - lane_pad * 2
            card_h = max(card_min_h, (avail_h - card_gap * (rows - 1)) / rows if rows > 0 else 0)
            
            for i, item_text in enumerate(items):
                card_top = area['top'] + lane_title_h + lane_pad + i * (card_h + card_gap)
                card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left + lane_pad, card_top, lane_w - 2 * lane_pad, card_h)
                card.fill.solid()
                card.fill.fore_color.rgb = RGBColor.from_string(CONFIG['COLORS']['card_bg'])
                card.line.color.rgb = RGBColor.from_string(CONFIG['COLORS']['card_border'])
                card.text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
                card.text_frame.margin_left = card.text_frame.margin_right = layout.px_to_emu(8)
                card.text_frame.word_wrap = True
                self._set_styled_text(card.text_frame, item_text, {'size': CONFIG['FONTS']['sizes']['body'], 'align': PP_ALIGN.CENTER})
                card_boxes[j].append({'left': left + lane_pad, 'top': card_top, 'width': lane_w - 2 * lane_pad, 'height': card_h})
        
        max_rows = max(len(cb) for cb in card_boxes) if card_boxes else 0
        for i in range(max_rows):
            for j in range(n - 1):
                if i < len(card_boxes[j]) and i < len(card_boxes[j+1]):
                    rect_a, rect_b = card_boxes[j][i], card_boxes[j+1][i]
                    from_x, to_x = rect_a['left'] + rect_a['width'], rect_b['left']
                    y_mid = rect_a['top'] + rect_a['height'] / 2
                    arrow_left = from_x + arrow_gap
                    arrow_width = to_x - from_x - 2 * arrow_gap
                    if arrow_width > 0:
                        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, arrow_left, y_mid - arrow_h/2, arrow_width, arrow_h)
                        arrow.fill.solid()
                        arrow.fill.fore_color.rgb = RGBColor.from_string(CONFIG['COLORS']['primary_blue'])
                        arrow.line.fill.background()

        self._draw_bottom_bar(slide, layout)
        self._add_google_footer(slide, layout, page_num)

    def _create_progress_slide(self, slide, data, layout, page_num):
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor.from_string(CONFIG['COLORS']['background_white'])
        self._draw_standard_title_header(slide, layout, 'contentSlide', data.get('title'))
        dy = self._draw_subhead_if_any(slide, layout, 'contentSlide', data.get('subhead'))
        
        area = layout.get_rect('contentSlide.body')
        area['top'] += dy
        items = data.get('items', [])
        n = len(items)
        if n == 0: return

        row_h = area['height'] / n
        for i, item in enumerate(items):
            y = area['top'] + i * row_h + row_h/2 - layout.px_to_emu(9)
            
            label_box = slide.shapes.add_textbox(area['left'], y, layout.px_to_emu(150), layout.px_to_emu(18))
            self._set_styled_text(label_box.text_frame, item.get('label', ''), {})

            bar_left = area['left'] + layout.px_to_emu(160)
            bar_w = area['width'] - layout.px_to_emu(220)
            bar_h = layout.px_to_emu(14)

            bg_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, bar_left, y, bar_w, bar_h)
            bg_bar.fill.solid()
            bg_bar.fill.fore_color.rgb = RGBColor.from_string(CONFIG['COLORS']['faint_gray'])
            bg_bar.line.fill.background()

            percent = max(0, min(100, item.get('percent', 0)))
            fg_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, bar_left, y, bar_w * (percent/100), bar_h)
            fg_bar.fill.solid()
            fg_bar.fill.fore_color.rgb = RGBColor.from_string(CONFIG['COLORS']['google_green'])
            fg_bar.line.fill.background()

            pct_box = slide.shapes.add_textbox(bar_left + bar_w + layout.px_to_emu(6), y - layout.px_to_emu(1), layout.px_to_emu(50), layout.px_to_emu(16))
            self._set_styled_text(pct_box.text_frame, f"{percent}%", {'size': CONFIG['FONTS']['sizes']['timelineLabel'], 'color': CONFIG['COLORS']['neutral_gray']})

        self._draw_bottom_bar(slide, layout)
        self._add_google_footer(slide, layout, page_num)

    def _create_closing_slide(self, slide, data, layout, page_num):
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor.from_string(CONFIG['COLORS']['background_white'])
        try:
            response = requests.get(CONFIG['LOGOS']['closing'])
            logo_data = io.BytesIO(response.content)
            with Image.open(logo_data) as img:
                aspect_ratio = img.height / img.width
            logo_w_emu = layout.px_to_emu(450)
            logo_h_emu = int(logo_w_emu * aspect_ratio)
            left, top = (layout.page_w_emu - logo_w_emu) / 2, (layout.page_h_emu - logo_h_emu) / 2
            logo_data.seek(0)
            slide.shapes.add_picture(logo_data, left, top, width=logo_w_emu)
        except Exception: pass

def generate_presentation(data):
    prs = Presentation()
    prs.slide_width, prs.slide_height = CONFIG['BASE_EMU']['W'], CONFIG['BASE_EMU']['H']
    layout_manager = LayoutManager(prs.slide_width, prs.slide_height)
    generator = SlideGenerator()
    page_counter = 0
    for item in data:
        if item['type'] not in ['title', 'closing']: page_counter += 1
        generator_func = generator.slide_generators.get(item['type'])
        if generator_func:
            slide = prs.slides.add_slide(prs.slide_layouts[6]) # 6 is a blank layout
            generator_func(slide, item, layout_manager, page_counter)
            if item.get('notes'):
                slide.notes_slide.notes_text_frame.text = item['notes']
    output_path = Path(SETTINGS['OUTPUT_FILENAME'])
    prs.save(str(output_path))
    print(f"プレゼンテーションを '{output_path.resolve()}' として保存しました。")

# ==============================================================================
# 4. メイン実行ブロック
# ==============================================================================
if __name__ == '__main__':
    try:
        generate_presentation(slide_data)
    except Exception as e:
        print(f"エラーが発生しました: {e}")
        import traceback
        traceback.print_exc()