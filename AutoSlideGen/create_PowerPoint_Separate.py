# -*- coding: utf-8 -*-

# ==============================================================================
# 必要なライブラリをインポートします。
# re: 正規表現を使って文字列から特定のパターンを抽出するために使用します。
# ast: 文字列を安全にPythonのデータ構造（リストや辞書など）に変換するために使用します。
# os, io, requests, datetime: ファイル操作、画像データの取り扱い、Webからのデータ取得、日付の扱いに使用します。
# pptx, PIL: PowerPointファイルの作成と画像サイズの計算に不可欠なライブラリです。
# ==============================================================================
import re
import ast
import io
import sys
import json
import requests
from pathlib import Path
from datetime import date
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image


# ==============================================================================
# 7.0 PYTHON_TEMPLATE_BLUEPRINT — 【Universal Python Design Ver. 2.0】完成済み設計図
#
# このセクションは、プロンプトVer. 2.0に含まれる固定のテンプレートコードです。
# 新しく追加されたスライドタイプを含む、全ての`slide_data`を解釈し、
# 実際にPowerPointスライドを組み立てる役割を持ちます。
# ==============================================================================

# --- 7.1. マスターデザイン設定 ---
SETTINGS = {
    'OUTPUT_FILENAME': 'presentation_extended_Separate.pptx'
}
CONFIG = {
    'BASE_PX': {'W': 960, 'H': 540},
    'BASE_EMU': {'W': Inches(10).emu, 'H': Inches(5.625).emu},
    'POS_PX': {
        'titleSlide': {
            'logo':      {'left': 55,  'top': 105, 'width': 135},
            'title':     {'left': 50,  'top': 230, 'width': 800, 'height': 90},
            'date':      {'left': 50,  'top': 340, 'width': 250, 'height': 40},
        },
        'contentSlide': {
            'headerLogo':    {'right': 20, 'top': 20, 'width': 75},
            'title':         {'left': 25, 'top': 60,  'width': 830, 'height': 65},
            'titleUnderline':{'left': 25, 'top': 128, 'width': 260, 'height': 4},
            'subhead':       {'left': 25, 'top': 140, 'width': 830, 'height': 30},
            'body':          {'left': 25, 'top': 172, 'width': 910, 'height': 303},
            'gridArea':      {'left': 25, 'top': 172, 'width': 910, 'height': 303},
            'compareLeft':   {'left': 25, 'top': 172, 'width': 430, 'height': 303},
            'compareRight':  {'left': 505, 'top': 172, 'width': 430, 'height': 303},
        },
        'sectionSlide': {
            'title':     {'left': 55, 'top': 230, 'width': 840, 'height': 80},
            'ghostNum':  {'left': 35, 'top': 120, 'width': 300, 'height': 200},
        },
        'footer': {
            'leftText':  {'left': 15, 'top': 505, 'width': 250, 'height': 20},
            'rightPage': {'right': 15, 'top': 505, 'width': 50,  'height': 20},
        },
        'bottomBar': {'left': 0, 'top': 534, 'width': 960, 'height': 6}
    },
    'FONTS': {
        'family': 'Arial',
        'sizes': {
            'title': 45, 'date': 16, 'sectionTitle': 38, 'contentTitle': 28,
            'subhead': 18, 'body': 14, 'footer': 9, 'cardTitle': 16,
            'cardDesc': 12, 'ghostNum': 180, 'processStep': 14, 'timelineLabel': 10,
            'laneTitle': 13,
        }
    },
    'COLORS': {
        'primary_blue': '4285F4', 'google_red': 'EA4335', 'google_yellow': 'FBBC04',
        'google_green': '34A853', 'text_primary': '333333', 'text_white': 'FFFFFF',
        'background_white': 'FFFFFF', 'background_gray': 'F8F9FA', 'card_bg': 'FFFFFF',
        'card_border': 'DADCE0', 'ghost_gray': 'EFEFED', 'faint_gray': 'E8EAED',
        'neutral_gray': '9E9E9E', 'lane_title_bg': 'F5F5F3', 'lane_border': 'DADCE0',
    },
    'LOGOS': {
        'header': 'https://upload.wikimedia.org/wikipedia/commons/thumb/2/2f/Google_2015_logo.svg/1024px-Google_2015_logo.svg.png',
        'closing': 'https://upload.wikimedia.org/wikipedia/commons/thumb/2/2f/Google_2015_logo.svg/1024px-Google_2015_logo.svg.png'
    },
    'FOOTER_TEXT': f"© {date.today().year} Your Organization"
}

# --- 7.2. コアロジック (PowerPoint生成の本体) ---
class LayoutManager:
    def __init__(self, page_w_emu, page_h_emu):
        self.page_w_emu, self.page_h_emu = page_w_emu, page_h_emu
        self.base_w_px, self.base_h_px = CONFIG['BASE_PX']['W'], CONFIG['BASE_PX']['H']
        self.scale_x, self.scale_y = self.page_w_emu / self.base_w_px, self.page_h_emu / self.base_h_px
    def px_to_emu(self, px, axis='x'):
        return int(px * (self.scale_x if axis == 'x' else self.scale_y))
    def get_rect(self, spec_path):
        keys = spec_path.split('.')
        pos_px = CONFIG['POS_PX']
        for key in keys: pos_px = pos_px[key]
        left_px = pos_px.get('left')
        if pos_px.get('right') is not None and left_px is None:
            left_px = self.base_w_px - pos_px['right'] - pos_px['width']
        return {
            'left': self.px_to_emu(left_px, 'x') if left_px is not None else None,
            'top': self.px_to_emu(pos_px['top'], 'y') if pos_px.get('top') is not None else None,
            'width': self.px_to_emu(pos_px['width'], 'x') if pos_px.get('width') is not None else None,
            'height': self.px_to_emu(pos_px['height'], 'y') if pos_px.get('height') is not None else None,
        }

class SlideGenerator:
    def __init__(self):
        self._section_counter = 0
        self.slide_generators = {
            'title': self._create_title_slide, 'section': self._create_section_slide,
            'content': self._create_content_slide, 'cards': self._create_cards_slide,
            'table': self._create_table_slide, 'closing': self._create_closing_slide,
            'compare': self._create_compare_slide, 'process': self._create_process_slide,
            'timeline': self._create_timeline_slide, 'diagram': self._create_diagram_slide,
            'progress': self._create_progress_slide,
        }
    def _set_font_style(self, run, style_opts):
        font = run.font
        font.name = style_opts.get('family', CONFIG['FONTS']['family'])
        font.size = Pt(style_opts.get('size', CONFIG['FONTS']['sizes']['body']))
        font.bold = style_opts.get('bold', False)
        color_str = style_opts.get('color', CONFIG['COLORS']['text_primary'])
        font.color.rgb = RGBColor.from_string(color_str)
    def _parse_inline_styles(self, text):
        pattern = r'(\*\*|\[\[)(.*?)\1'
        parts, last_end = [], 0
        for match in re.finditer(pattern, text):
            start, end = match.span()
            tag, content = match.groups()
            if start > last_end: parts.append({'text': text[last_end:start], 'style': {}})
            style = {'bold': True, 'color': CONFIG['COLORS']['primary_blue']} if tag == '[[' else {'bold': True}
            parts.append({'text': content, 'style': style})
            last_end = end
        if last_end < len(text): parts.append({'text': text[last_end:], 'style': {}})
        return parts
    def _set_styled_text(self, text_frame, text, base_style_opts):
        text_frame.clear()
        p = text_frame.paragraphs[0]
        p.font.name = base_style_opts.get('family', CONFIG['FONTS']['family'])
        p.font.size = Pt(base_style_opts.get('size', CONFIG['FONTS']['sizes']['body']))
        p.alignment = base_style_opts.get('align', PP_ALIGN.LEFT)
        for i, line in enumerate(text.split('\n')):
            if i > 0: p = text_frame.add_paragraph()
            parts = self._parse_inline_styles(line) or [{'text': line, 'style': {}}]
            for part in parts:
                run = p.add_run()
                run.text = part['text']
                self._set_font_style(run, {**base_style_opts, **part['style']})
    def _set_bullets_with_inline_styles(self, text_frame, points, base_style_opts):
        text_frame.clear()
        for i, point_text in enumerate(points):
            p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
            run_bullet = p.add_run()
            run_bullet.text = "• "
            self._set_font_style(run_bullet, base_style_opts)
            parts = self._parse_inline_styles(point_text)
            for part in parts:
                run = p.add_run()
                run.text = part['text']
                self._set_font_style(run, {**base_style_opts, **part['style']})
    def _draw_bottom_bar(self, slide, layout):
        rect = layout.get_rect('bottomBar')
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, rect['left'], rect['top'], rect['width'], rect['height'])
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor.from_string(CONFIG['COLORS']['primary_blue'])
        shape.line.fill.background()
    def _add_google_footer(self, slide, layout, page_num):
        rect = layout.get_rect('footer.leftText')
        textbox = slide.shapes.add_textbox(rect['left'], rect['top'], rect['width'], rect['height'])
        self._set_styled_text(textbox.text_frame, CONFIG['FOOTER_TEXT'], {'size': CONFIG['FONTS']['sizes']['footer']})
        if page_num > 0:
            rect = layout.get_rect('footer.rightPage')
            textbox = slide.shapes.add_textbox(rect['left'], rect['top'], rect['width'], rect['height'])
            self._set_styled_text(textbox.text_frame, str(page_num), {
                'size': CONFIG['FONTS']['sizes']['footer'], 'color': CONFIG['COLORS']['primary_blue'], 'align': PP_ALIGN.RIGHT
            })
    def _draw_standard_title_header(self, slide, layout, key, title):
        rect = layout.get_rect(f'{key}.headerLogo')
        try:
            response = requests.get(CONFIG['LOGOS']['header'])
            slide.shapes.add_picture(io.BytesIO(response.content), rect['left'], rect['top'], width=rect['width'])
        except Exception: pass
        rect = layout.get_rect(f'{key}.title')
        textbox = slide.shapes.add_textbox(rect['left'], rect['top'], rect['width'], rect['height'])
        self._set_styled_text(textbox.text_frame, title, {'size': CONFIG['FONTS']['sizes']['contentTitle'], 'bold': True})
        rect = layout.get_rect(f'{key}.titleUnderline')
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, rect['left'], rect['top'], rect['width'], rect['height'])
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor.from_string(CONFIG['COLORS']['primary_blue'])
        shape.line.fill.background()
    def _draw_subhead_if_any(self, slide, layout, key, subhead):
        if not subhead: return 0
        rect = layout.get_rect(f'{key}.subhead')
        textbox = slide.shapes.add_textbox(rect['left'], rect['top'], rect['width'], rect['height'])
        self._set_styled_text(textbox.text_frame, subhead, {'size': CONFIG['FONTS']['sizes']['subhead']})
        return layout.px_to_emu(36)
    def _create_title_slide(self, slide, data, layout, page_num):
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor.from_string(CONFIG['COLORS']['background_white'])
        rect = layout.get_rect('titleSlide.logo')
        try:
            response = requests.get(CONFIG['LOGOS']['header'])
            slide.shapes.add_picture(io.BytesIO(response.content), rect['left'], rect['top'], width=rect['width'])
        except Exception: pass
        rect = layout.get_rect('titleSlide.title')
        textbox = slide.shapes.add_textbox(rect['left'], rect['top'], rect['width'], rect['height'])
        self._set_styled_text(textbox.text_frame, data.get('title', ''), {'size': CONFIG['FONTS']['sizes']['title'], 'bold': True})
        rect = layout.get_rect('titleSlide.date')
        textbox = slide.shapes.add_textbox(rect['left'], rect['top'], rect['width'], rect['height'])
        self._set_styled_text(textbox.text_frame, data.get('date', ''), {'size': CONFIG['FONTS']['sizes']['date']})
        self._draw_bottom_bar(slide, layout)
    def _create_section_slide(self, slide, data, layout, page_num):
        self._section_counter += 1
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor.from_string(CONFIG['COLORS']['background_gray'])
        num_str = str(data.get('sectionNo', self._section_counter)).zfill(2)
        rect = layout.get_rect('sectionSlide.ghostNum')
        textbox = slide.shapes.add_textbox(rect['left'], rect['top'], rect['width'], rect['height'])
        self._set_styled_text(textbox.text_frame, num_str, {
            'size': CONFIG['FONTS']['sizes']['ghostNum'], 'bold': True, 'color': CONFIG['COLORS']['ghost_gray'], 'align': PP_ALIGN.CENTER
        })
        rect = layout.get_rect('sectionSlide.title')
        textbox = slide.shapes.add_textbox(rect['left'], rect['top'], rect['width'], rect['height'])
        self._set_styled_text(textbox.text_frame, data.get('title', ''), {
            'size': CONFIG['FONTS']['sizes']['sectionTitle'], 'bold': True, 'align': PP_ALIGN.CENTER
        })
        self._add_google_footer(slide, layout, page_num)
    def _create_content_slide(self, slide, data, layout, page_num):
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor.from_string(CONFIG['COLORS']['background_white'])
        self._draw_standard_title_header(slide, layout, 'contentSlide', data.get('title'))
        dy = self._draw_subhead_if_any(slide, layout, 'contentSlide', data.get('subhead'))
        body_rect = layout.get_rect('contentSlide.body')
        body_rect['top'] += dy
        textbox = slide.shapes.add_textbox(body_rect['left'], body_rect['top'], body_rect['width'], body_rect['height'])
        self._set_bullets_with_inline_styles(textbox.text_frame, data.get('points', []), {})
        self._draw_bottom_bar(slide, layout)
        self._add_google_footer(slide, layout, page_num)
    def _create_cards_slide(self, slide, data, layout, page_num):
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor.from_string(CONFIG['COLORS']['background_white'])
        self._draw_standard_title_header(slide, layout, 'contentSlide', data.get('title'))
        dy = self._draw_subhead_if_any(slide, layout, 'contentSlide', data.get('subhead'))
        area = layout.get_rect('contentSlide.gridArea')
        area['top'] += dy
        items = data.get('items', [])
        cols = data.get('columns', 3)
        rows = -(-len(items) // cols)
        gap = layout.px_to_emu(16)
        card_w = (area['width'] - gap * (cols - 1)) // cols
        card_h = (area['height'] - gap * (rows - 1)) // rows
        for i, item in enumerate(items):
            r, c = divmod(i, cols)
            left, top = area['left'] + c * (card_w + gap), area['top'] + r * (card_h + gap)
            shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, card_w, card_h)
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor.from_string(CONFIG['COLORS']['card_bg'])
            shape.line.color.rgb = RGBColor.from_string(CONFIG['COLORS']['card_border'])
            shape.line.width = Pt(1).emu
            tf = shape.text_frame
            tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = layout.px_to_emu(10)
            tf.word_wrap = True
            tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
            if isinstance(item, dict):
                p_title = tf.paragraphs[0]
                run_title = p_title.add_run()
                run_title.text = item.get('title', '')
                self._set_font_style(run_title, {'size': CONFIG['FONTS']['sizes']['cardTitle'], 'bold': True})
                
                if item.get('desc'):
                    p_desc = tf.add_paragraph()
                    run_desc = p_desc.add_run()
                    run_desc.text = item.get('desc', '')
                    self._set_font_style(run_desc, {'size': CONFIG['FONTS']['sizes']['cardDesc']})
        self._draw_bottom_bar(slide, layout)
        self._add_google_footer(slide, layout, page_num)
    def _create_table_slide(self, slide, data, layout, page_num):
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor.from_string(CONFIG['COLORS']['background_white'])
        self._draw_standard_title_header(slide, layout, 'contentSlide', data.get('title'))
        dy = self._draw_subhead_if_any(slide, layout, 'contentSlide', data.get('subhead'))
        area = layout.get_rect('contentSlide.body')
        area['top'] += dy
        headers = data.get('headers', [])
        rows_data = data.get('rows', [])
        if not headers or not rows_data: return
        shape = slide.shapes.add_table(len(rows_data) + 1, len(headers), area['left'], area['top'], area['width'], area['height'])
        table = shape.table
        for c, header_text in enumerate(headers):
            cell = table.cell(0, c)
            cell.text = header_text
            cell.text_frame.paragraphs[0].font.bold = True
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        for r, row_data in enumerate(rows_data):
            for c, cell_text in enumerate(row_data):
                cell = table.cell(r + 1, c)
                cell.text = cell_text
                cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        self._draw_bottom_bar(slide, layout)
        self._add_google_footer(slide, layout, page_num)
    def _create_compare_slide(self, slide, data, layout, page_num):
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor.from_string(CONFIG['COLORS']['background_white'])
        self._draw_standard_title_header(slide, layout, 'contentSlide', data.get('title'))
        dy = self._draw_subhead_if_any(slide, layout, 'contentSlide', data.get('subhead'))
        def draw_compare_box(rect, title, items):
            box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, rect['left'], rect['top'], rect['width'], rect['height'])
            box.fill.solid()
            box.fill.fore_color.rgb = RGBColor.from_string(CONFIG['COLORS']['background_gray'])
            box.line.color.rgb = RGBColor.from_string(CONFIG['COLORS']['card_border'])
            box.line.width = Pt(1).emu
            title_h = layout.px_to_emu(40)
            title_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, rect['left'], rect['top'], rect['width'], title_h)
            title_bar.fill.solid()
            title_bar.fill.fore_color.rgb = RGBColor.from_string(CONFIG['COLORS']['primary_blue'])
            title_bar.line.fill.background()
            title_bar.text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
            self._set_styled_text(title_bar.text_frame, title, {'size': CONFIG['FONTS']['sizes']['laneTitle'], 'bold': True, 'color': CONFIG['COLORS']['text_white'], 'align': PP_ALIGN.CENTER})
            pad = layout.px_to_emu(12)
            text_box = slide.shapes.add_textbox(rect['left'] + pad, rect['top'] + title_h + pad, rect['width'] - 2*pad, rect['height'] - title_h - 2*pad)
            self._set_bullets_with_inline_styles(text_box.text_frame, items, {})
        left_rect = layout.get_rect('contentSlide.compareLeft')
        left_rect['top'] += dy
        draw_compare_box(left_rect, data.get('leftTitle', ''), data.get('leftItems', []))
        right_rect = layout.get_rect('contentSlide.compareRight')
        right_rect['top'] += dy
        draw_compare_box(right_rect, data.get('rightTitle', ''), data.get('rightItems', []))
        self._draw_bottom_bar(slide, layout)
        self._add_google_footer(slide, layout, page_num)
    def _create_process_slide(self, slide, data, layout, page_num):
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor.from_string(CONFIG['COLORS']['background_white'])
        self._draw_standard_title_header(slide, layout, 'contentSlide', data.get('title'))
        dy = self._draw_subhead_if_any(slide, layout, 'contentSlide', data.get('subhead'))
        area = layout.get_rect('contentSlide.body')
        area['top'] += dy
        steps = data.get('steps', [])
        n = len(steps)
        if n == 0: return
        box_h = layout.px_to_emu(40)
        gap_y = (area['height'] - n * box_h) / max(1, n - 1) if n > 1 else 0
        arrow_w, arrow_h = layout.px_to_emu(20), layout.px_to_emu(30)
        for i, step_text in enumerate(steps):
            y = area['top'] + i * (box_h + gap_y)
            num_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, area['left'], y, box_h, box_h)
            num_box.fill.solid()
            num_box.fill.fore_color.rgb = RGBColor.from_string(CONFIG['COLORS']['primary_blue'])
            num_box.line.fill.background()
            num_box.text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
            self._set_styled_text(num_box.text_frame, str(i + 1), {'bold': True, 'color': CONFIG['COLORS']['text_white'], 'align': PP_ALIGN.CENTER})
            txt_box = slide.shapes.add_textbox(area['left'] + box_h + layout.px_to_emu(15), y, area['width'] - box_h - layout.px_to_emu(15), box_h)
            txt_box.text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
            self._set_styled_text(txt_box.text_frame, step_text, {'size': CONFIG['FONTS']['sizes']['processStep']})
            if i < n - 1:
                arrow_y = y + box_h + (gap_y - arrow_h) / 2
                arrow = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, area['left'] + (box_h - arrow_w)/2, arrow_y, arrow_w, arrow_h)
                arrow.fill.solid()
                arrow.fill.fore_color.rgb = RGBColor.from_string(CONFIG['COLORS']['faint_gray'])
                arrow.line.fill.background()
        self._draw_bottom_bar(slide, layout)
        self._add_google_footer(slide, layout, page_num)
    def _create_timeline_slide(self, slide, data, layout, page_num):
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor.from_string(CONFIG['COLORS']['background_white'])
        self._draw_standard_title_header(slide, layout, 'contentSlide', data.get('title'))
        dy = self._draw_subhead_if_any(slide, layout, 'contentSlide', data.get('subhead'))
        area = layout.get_rect('contentSlide.body')
        area['top'] += dy
        milestones = data.get('milestones', [])
        n = len(milestones)
        if n == 0: return
        inner_margin = layout.px_to_emu(60)
        base_y = area['top'] + area['height'] * 0.5
        left_x, right_x = area['left'] + inner_margin, area['left'] + area['width'] - inner_margin
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left_x, base_y - layout.px_to_emu(1), right_x - left_x, layout.px_to_emu(2))
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor.from_string(CONFIG['COLORS']['faint_gray'])
        line.line.fill.background()
        dot_r = layout.px_to_emu(12)
        gap_x = (right_x - left_x) / max(1, n - 1) if n > 1 else 0
        for i, m in enumerate(milestones):
            x = left_x + gap_x * i
            dot = slide.shapes.add_shape(MSO_SHAPE.ELLIPSE, x - dot_r/2, base_y - dot_r/2, dot_r, dot_r)
            state_colors = {'done': CONFIG['COLORS']['google_green'], 'next': CONFIG['COLORS']['google_yellow'], 'todo': CONFIG['COLORS']['neutral_gray']}
            state = m.get('state', 'todo').lower()
            dot.fill.solid()
            dot.fill.fore_color.rgb = RGBColor.from_string(state_colors.get(state, state_colors['todo']))
            dot.line.fill.background()
            label_box = slide.shapes.add_textbox(x - layout.px_to_emu(50), base_y - layout.px_to_emu(50), layout.px_to_emu(100), layout.px_to_emu(30))
            self._set_styled_text(label_box.text_frame, m.get('label', ''), {'size': CONFIG['FONTS']['sizes']['timelineLabel'], 'bold': True, 'align': PP_ALIGN.CENTER})
            date_box = slide.shapes.add_textbox(x - layout.px_to_emu(50), base_y + layout.px_to_emu(15), layout.px_to_emu(100), layout.px_to_emu(20))
            self._set_styled_text(date_box.text_frame, m.get('date', ''), {'size': CONFIG['FONTS']['sizes']['timelineLabel'], 'color': CONFIG['COLORS']['neutral_gray'], 'align': PP_ALIGN.CENTER})
        self._draw_bottom_bar(slide, layout)
        self._add_google_footer(slide, layout, page_num)
    def _create_diagram_slide(self, slide, data, layout, page_num):
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor.from_string(CONFIG['COLORS']['background_white'])
        self._draw_standard_title_header(slide, layout, 'contentSlide', data.get('title'))
        dy = self._draw_subhead_if_any(slide, layout, 'contentSlide', data.get('subhead'))
        area = layout.get_rect('contentSlide.body')
        area['top'] += dy
        lanes = data.get('lanes', [])
        n = len(lanes)
        if n == 0: return
        lane_gap = layout.px_to_emu(24)
        lane_pad = layout.px_to_emu(10)
        lane_title_h = layout.px_to_emu(30)
        card_gap = layout.px_to_emu(12)
        card_min_h = layout.px_to_emu(48)
        arrow_h, arrow_gap = layout.px_to_emu(10), layout.px_to_emu(8)
        lane_w = (area['width'] - lane_gap * (n - 1)) / n
        card_boxes = [[] for _ in range(n)]
        for j, lane in enumerate(lanes):
            left = area['left'] + j * (lane_w + lane_gap)
            lt = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, area['top'], lane_w, lane_title_h)
            lt.fill.solid()
            lt.fill.fore_color.rgb = RGBColor.from_string(CONFIG['COLORS']['lane_title_bg'])
            lt.line.color.rgb = RGBColor.from_string(CONFIG['COLORS']['lane_border'])
            lt.text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
            self._set_styled_text(lt.text_frame, lane.get('title', ''), {'size': CONFIG['FONTS']['sizes']['laneTitle'], 'bold': True, 'align': PP_ALIGN.CENTER})
            items = lane.get('items', [])
            rows = len(items)
            avail_h = area['height'] - lane_title_h - lane_pad * 2
            card_h = max(card_min_h, (avail_h - card_gap * (rows - 1)) / rows if rows > 0 else 0)
            for i, item_text in enumerate(items):
                card_top = area['top'] + lane_title_h + lane_pad + i * (card_h + card_gap)
                card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left + lane_pad, card_top, lane_w - 2 * lane_pad, card_h)
                card.fill.solid()
                card.fill.fore_color.rgb = RGBColor.from_string(CONFIG['COLORS']['card_bg'])
                card.line.color.rgb = RGBColor.from_string(CONFIG['COLORS']['card_border'])
                card.text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
                card.text_frame.margin_left = card.text_frame.margin_right = layout.px_to_emu(8)
                card.text_frame.word_wrap = True
                self._set_styled_text(card.text_frame, item_text, {'size': CONFIG['FONTS']['sizes']['body'], 'align': PP_ALIGN.CENTER})
                card_boxes[j].append({'left': left + lane_pad, 'top': card_top, 'width': lane_w - 2 * lane_pad, 'height': card_h})
        max_rows = max(len(cb) for cb in card_boxes) if card_boxes else 0
        for i in range(max_rows):
            for j in range(n - 1):
                if i < len(card_boxes[j]) and i < len(card_boxes[j+1]):
                    rect_a, rect_b = card_boxes[j][i], card_boxes[j+1][i]
                    from_x, to_x = rect_a['left'] + rect_a['width'], rect_b['left']
                    y_mid = rect_a['top'] + rect_a['height'] / 2
                    arrow_left, arrow_width = from_x + arrow_gap, to_x - from_x - 2 * arrow_gap
                    if arrow_width > 0:
                        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, arrow_left, y_mid - arrow_h/2, arrow_width, arrow_h)
                        arrow.fill.solid()
                        arrow.fill.fore_color.rgb = RGBColor.from_string(CONFIG['COLORS']['primary_blue'])
                        arrow.line.fill.background()
        self._draw_bottom_bar(slide, layout)
        self._add_google_footer(slide, layout, page_num)
    def _create_progress_slide(self, slide, data, layout, page_num):
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor.from_string(CONFIG['COLORS']['background_white'])
        self._draw_standard_title_header(slide, layout, 'contentSlide', data.get('title'))
        dy = self._draw_subhead_if_any(slide, layout, 'contentSlide', data.get('subhead'))
        area = layout.get_rect('contentSlide.body')
        area['top'] += dy
        items = data.get('items', [])
        n = len(items)
        if n == 0: return
        row_h = area['height'] / n
        for i, item in enumerate(items):
            y = area['top'] + i * row_h + row_h/2 - layout.px_to_emu(9)
            label_box = slide.shapes.add_textbox(area['left'], y, layout.px_to_emu(150), layout.px_to_emu(18))
            self._set_styled_text(label_box.text_frame, item.get('label', ''), {})
            bar_left, bar_w, bar_h = area['left'] + layout.px_to_emu(160), area['width'] - layout.px_to_emu(220), layout.px_to_emu(14)
            bg_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, bar_left, y, bar_w, bar_h)
            bg_bar.fill.solid()
            bg_bar.fill.fore_color.rgb = RGBColor.from_string(CONFIG['COLORS']['faint_gray'])
            bg_bar.line.fill.background()
            percent = max(0, min(100, item.get('percent', 0)))
            fg_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, bar_left, y, bar_w * (percent/100), bar_h)
            fg_bar.fill.solid()
            fg_bar.fill.fore_color.rgb = RGBColor.from_string(CONFIG['COLORS']['google_green'])
            fg_bar.line.fill.background()
            pct_box = slide.shapes.add_textbox(bar_left + bar_w + layout.px_to_emu(6), y - layout.px_to_emu(1), layout.px_to_emu(50), layout.px_to_emu(16))
            self._set_styled_text(pct_box.text_frame, f"{percent}%", {'size': CONFIG['FONTS']['sizes']['timelineLabel'], 'color': CONFIG['COLORS']['neutral_gray']})
        self._draw_bottom_bar(slide, layout)
        self._add_google_footer(slide, layout, page_num)
    def _create_closing_slide(self, slide, data, layout, page_num):
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor.from_string(CONFIG['COLORS']['background_white'])
        try:
            response = requests.get(CONFIG['LOGOS']['closing'])
            logo_data = io.BytesIO(response.content)
            with Image.open(logo_data) as img:
                aspect_ratio = img.height / img.width
            logo_w_emu = layout.px_to_emu(450)
            logo_h_emu = int(logo_w_emu * aspect_ratio)
            left, top = (layout.page_w_emu - logo_w_emu) / 2, (layout.page_h_emu - logo_h_emu) / 2
            logo_data.seek(0)
            slide.shapes.add_picture(logo_data, left, top, width=logo_w_emu)
        except Exception: pass

def generate_presentation(data, output_dir='/tmp'):
    prs = Presentation()
    prs.slide_width, prs.slide_height = CONFIG['BASE_EMU']['W'], CONFIG['BASE_EMU']['H']
    layout_manager = LayoutManager(prs.slide_width, prs.slide_height)
    generator = SlideGenerator()
    page_counter = 0
    for item in data:
        if item['type'] not in ['title', 'closing']: page_counter += 1
        generator_func = generator.slide_generators.get(item['type'])
        if generator_func:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            generator_func(slide, item, layout_manager, page_counter)
            if item.get('notes'):
                slide.notes_slide.notes_text_frame.text = item['notes']
    output_path = Path(output_dir) / SETTINGS['OUTPUT_FILENAME']
    prs.save(str(output_path))
    print(f"プレゼンテーションを '{output_path}' として保存しました。")
    return str(output_path)

# ==============================================================================
# 安全な変換と実行を行うメインの処理関数
# ==============================================================================
def safe_slide_generator_from_ai_script(slide_data_string: str):
    """
    AIが生成したslide_dataのリスト文字列を直接引数として受け取り、
    スライド生成を実行する。
    
    Args:
        slide_data_string (str): AIによって生成された、リスト形式のPython文字列。 例: "[{'type': 'title', 'title': 'サンプル'}]"

    Returns:
        str or None: 成功した場合は生成されたPowerPointファイルのパス、失敗した場合はNoneを返す。
    """
    print("--- AIからの出力文字列の処理を開始します ---")

    # 受け取った文字列の前後の空白を削除
    processed_string = slide_data_string.strip()

    # 引数がリスト形式の文字列（'['で始まり']'で終わる）か簡易的にチェック
    if not (processed_string.startswith('[') and processed_string.endswith(']')):
        print("エラー: 入力された文字列がリスト形式ではありません。'['で始まり']'で終わる必要があります。")
        return None

    try:
        # 文字列を安全にPythonのリストオブジェクトに変換
        slide_data_list = ast.literal_eval(processed_string)
        print("文字列を安全なPythonリストオブジェクトに変換しました。")

    except (ValueError, SyntaxError) as e:
        print(f"エラー: 文字列の形式が不正です。変換に失敗しました。詳細: {e}")
        return None

    # 変換後のデータがリスト形式であり、中身が空でないことを確認
    if isinstance(slide_data_list, list) and slide_data_list:
        print(f"{len(slide_data_list)}スライド分のデータを元に、PowerPointの生成を開始します。")
        # スライド生成関数を呼び出し
        generated_file_path = generate_presentation(slide_data_list)
        return generated_file_path
    else:
        print("エラー: 変換後のデータが空、またはリスト形式ではありません。処理を中断します。")
        return None

# ==============================================================================
# 実行メインブロック
# ==============================================================================
if __name__ == '__main__':
    if len(sys.argv) > 1:
        # コマンドライン引数の2番目（[1]）をデータとして受け取る
        ai_output_string = sys.argv[1]

        # /tmp ディレクトリが存在しない場合は作成する
        tmp_dir = Path('/tmp')
        if not tmp_dir.exists():
            tmp_dir.mkdir(parents=True)
            
        safe_slide_generator_from_ai_script(ai_output_string)
    else:
        print("エラー: コマンドライン引数としてslide_dataの文字列を指定してください。")



# ==============================================================================
# AWS Lambda ハンドラ関数
# ==============================================================================
def lambda_handler(event, context):
    """
    AWS Lambdaのエントリーポイント。
    Salesforceからのリクエスト(event)を処理する。
    """
    print("--- Lambda関数がトリガーされました ---")
    
    try:
        # Salesforceから渡されるデータは event['body'] にJSON形式で入っている
        # slide_data_string = json.loads(event['body'])['slideData'] # Salesforce側の実装に合わせてキーを調整
        
        # API Gatewayを使わないテスト用の単純なペイロードを想定
        slide_data_string = event['slideData']
        
        print("ペイロードからslide_data文字列の取得に成功しました。")
        
        # 既存のロジックを呼び出してPowerPointを生成
        # Lambdaでは /tmp ディレクトリのみ書き込み可能
        output_path = generate_presentation(ast.literal_eval(slide_data_string), output_dir='./')
        
        # S3へのアップロード処理などをここに追加するのが一般的
        # （今回はファイルパスを返すだけの簡易的な例）
        
        print(f"プレゼンテーションの生成に成功しました: {output_path}")
        
        # Salesforceに成功レスポンスを返す
        return {
            'statusCode': 200,
            'body': json.dumps({
                'message': 'プレゼンテーションが正常に作成されました。',
                'filePath': output_path
            })
        }

    except Exception as e:
        print(f"エラーが発生しました: {e}")
        # Salesforceにエラーレスポンスを返す
        return {
            'statusCode': 500,
            'body': json.dumps({'error': str(e)})
        }
