#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Lambda環境での依存関係インポートテスト
"""
import sys
import json
from pathlib import Path

def lambda_handler(event, context):
    """
    Lambda関数のテストハンドラー - 依存関係のインポートテスト
    """
    results = []
    errors = []
    
    try:
        # lxml.etreeのインポートテスト（最重要）
        from lxml import etree
        results.append("✅ lxml.etree インポート成功")
        
        # python-pptxのインポートテスト
        from pptx import Presentation
        results.append("✅ python-pptx インポート成功")
        
        # requestsのインポートテスト  
        import requests
        results.append("✅ requests インポート成功")
        
        # PILのインポートテスト
        from PIL import Image
        results.append("✅ PIL インポート成功")
        
        # 簡単な動作テスト
        # PowerPoint作成テスト
        prs = Presentation()
        slide_layout = prs.slide_layouts[0] 
        slide = prs.slides.add_slide(slide_layout)
        results.append("✅ PowerPoint作成テスト成功")
        
        # lxml動作テスト
        root = etree.Element("test")
        etree.SubElement(root, "child").text = "Hello"
        xml_str = etree.tostring(root, encoding='unicode')
        results.append("✅ lxml動作テスト成功")
        
        # HTTPリクエストテスト（モック）
        response = requests.get("https://httpbin.org/get", timeout=10)
        if response.status_code == 200:
            results.append("✅ HTTP通信テスト成功")
        
    except Exception as e:
        error_msg = f"❌ エラー: {str(e)}"
        errors.append(error_msg)
    
    # 結果をまとめる
    success_count = len(results)
    error_count = len(errors)
    
    response = {
        'statusCode': 200 if error_count == 0 else 500,
        'body': json.dumps({
            'message': f'依存関係テスト完了: 成功{success_count}件、エラー{error_count}件',
            'results': results,
            'errors': errors,
            'success': error_count == 0
        }, ensure_ascii=False)
    }
    
    return response

if __name__ == '__main__':
    # ローカルテスト実行
    print("🔧 Lambda依存関係テストを開始...")
    result = lambda_handler({}, {})
    print(json.dumps(json.loads(result['body']), indent=2, ensure_ascii=False))