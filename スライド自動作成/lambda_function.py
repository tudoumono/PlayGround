# -*- coding: utf-8 -*-

# ==============================================================================
# 必要なライブラリをインポートします
# ==============================================================================
import json
import os
import io
import re
import ast
import uuid
import boto3
from datetime import date
from botocore.exceptions import ClientError
import requests
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

# ==============================================================================
# AWS S3を操作するためのクライアントを初期化します
# ==============================================================================
s3_client = boto3.client('s3')

# ==============================================================================
# 1. マスターデザイン設定
#
# この`CONFIG`辞書は、生成されるPowerPointスライド全体のデザインを定義する「設計図」です。
# Lambdaの環境変数から動的に設定を読み込むことで、コードを変更せずにデザインを調整できます。
# ==============================================================================

# 環境変数からロゴURLやフォント名などを読み込みます。
# os.environ.get('キー', 'デフォルト値') を使うことで、環境変数が未設定の場合でもエラーを防ぎます。
logo_header_url = os.environ.get(
    'LOGO_HEADER_URL',
    'https://upload.wikimedia.org/wikipedia/commons/thumb/2/2f/Google_2015_logo.svg/1024px-Google_2015_logo.svg.png'
)
logo_closing_url = os.environ.get(
    'LOGO_CLOSING_URL',
    'https://upload.wikimedia.org/wikipedia/commons/thumb/2/2f/Google_2015_logo.svg/1024px-Google_2015_logo.svg.png'
)
footer_org_name = os.environ.get('FOOTER_ORGANIZATION_NAME', 'Your Organization')
default_font = os.environ.get('DEFAULT_FONT_FAMILY', 'Arial')
primary_color = os.environ.get('THEME_COLOR_PRIMARY', '4285F4')

# デザイン設定を格納するCONFIG辞書
CONFIG = {
    # スライドの基本サイズ (ピクセルとPowerPoint内部単位EMU)
    'BASE_PX': {'W': 960, 'H': 540},
    'BASE_EMU': {'W': Inches(10).emu, 'H': Inches(5.625).emu},

    # 各要素の配置とサイズ定義 (ピクセル単位)
    'POS_PX': {
        'titleSlide': {
            'logo':      {'left': 55,  'top': 105, 'width': 135},
            'title':     {'left': 50,  'top': 230, 'width': 800, 'height': 90},
            'date':      {'left': 50,  'top': 340, 'width': 250, 'height': 40},
        },
        'contentSlide': {
            'headerLogo':    {'right': 20, 'top': 20, 'width': 75},
            'title':         {'left': 25, 'top': 60,  'width': 830, 'height': 65},
            'titleUnderline':{'left': 25, 'top': 128, 'width': 260, 'height': 4},
            'subhead':       {'left': 25, 'top': 140, 'width': 830, 'height': 30},
            'body':          {'left': 25, 'top': 172, 'width': 910, 'height': 303},
            'gridArea':      {'left': 25, 'top': 172, 'width': 910, 'height': 303},
            'compareLeft':   {'left': 25, 'top': 172, 'width': 430, 'height': 303},
            'compareRight':  {'left': 505, 'top': 172, 'width': 430, 'height': 303},
        },
        'sectionSlide': {
            'title':    {'left': 55, 'top': 230, 'width': 840, 'height': 80},
            'ghostNum': {'left': 35, 'top': 120, 'width': 300, 'height': 200},
        },
        'footer': {
            'leftText':  {'left': 15, 'top': 505, 'width': 250, 'height': 20},
            'rightPage': {'right': 15, 'top': 505, 'width': 50,  'height': 20},
        },
        'bottomBar': {'left': 0, 'top': 534, 'width': 960, 'height': 6}
    },

    # フォント設定
    'FONTS': {
        'family': default_font,
        'sizes': {
            'title': 45, 'date': 16, 'sectionTitle': 38, 'contentTitle': 28,
            'subhead': 18, 'body': 14, 'footer': 9, 'cardTitle': 16,
            'cardDesc': 12, 'ghostNum': 180, 'processStep': 14, 'timelineLabel': 10,
            'laneTitle': 13,
        }
    },

    # カラーパレット設定
    'COLORS': {
        'primary_blue': primary_color,
        'google_red': 'EA4335', 'google_yellow': 'FBBC04',
        'google_green': '34A853', 'text_primary': '333333', 'text_white': 'FFFFFF',
        'background_white': 'FFFFFF', 'background_gray': 'F8F9FA', 'card_bg': 'FFFFFF',
        'card_border': 'DADCE0', 'ghost_gray': 'EFEFED', 'faint_gray': 'E8EAED',
        'neutral_gray': '9E9E9E', 'lane_title_bg': 'F5F5F3', 'lane_border': 'DADCE0',
    },

    # ロゴ画像の設定
    'LOGOS': {
        'header': logo_header_url,
        'closing': logo_closing_url
    },

    # フッターテキストの設定
    'FOOTER_TEXT': f"© {date.today().year} {footer_org_name}"
}

# ==============================================================================
# 2. コアロジック (PowerPoint生成の本体)
# ==============================================================================
class LayoutManager:
    """
    ピクセル単位で定義されたデザイン設定を、PowerPointの内部単位(EMU)に変換し、
    各要素の正確な座標とサイズを計算する責務を持つクラス。
    """
    def __init__(self, page_w_emu, page_h_emu):
        self.page_w_emu, self.page_h_emu = page_w_emu, page_h_emu
        self.base_w_px, self.base_h_px = CONFIG['BASE_PX']['W'], CONFIG['BASE_PX']['H']
        # ピクセルからEMUへの変換スケールを計算
        self.scale_x, self.scale_y = self.page_w_emu / self.base_w_px, self.page_h_emu / self.base_h_px

    def px_to_emu(self, px, axis='x'):
        """ピクセル値をEMU値に変換する。"""
        return int(px * (self.scale_x if axis == 'x' else self.scale_y))

    def get_rect(self, spec_path):
        """CONFIG内のパス(例: 'titleSlide.logo')を指定して、EMU単位の矩形情報を取得する。"""
        keys = spec_path.split('.')
        pos_px = CONFIG['POS_PX']
        for key in keys: pos_px = pos_px[key]
        
        left_px = pos_px.get('left')
        # 'right'指定がある場合は、右端からの距離として左端座標を計算
        if pos_px.get('right') is not None and left_px is None:
            left_px = self.base_w_px - pos_px['right'] - pos_px['width']
            
        return {
            'left': self.px_to_emu(left_px, 'x') if left_px is not None else None,
            'top': self.px_to_emu(pos_px['top'], 'y') if pos_px.get('top') is not None else None,
            'width': self.px_to_emu(pos_px['width'], 'x') if pos_px.get('width') is not None else None,
            'height': self.px_to_emu(pos_px['height'], 'y') if pos_px.get('height') is not None else None,
        }

class SlideGenerator:
    """
    slide_dataの各要素に基づき、対応するスライドを実際に描画するメソッド群を持つクラス。
    PowerPointファイルの中身を作成する心臓部。
    """
    def __init__(self):
        self._section_counter = 0
        # slide_dataの'type'キーと、それを描画するメソッドを対応付ける辞書
        self.slide_generators = {
            'title': self._create_title_slide,
            'section': self._create_section_slide,
            'content': self._create_content_slide,
            'cards': self._create_cards_slide,
            'table': self._create_table_slide,
            'closing': self._create_closing_slide,
            'compare': self._create_compare_slide,
            'process': self._create_process_slide,
            'timeline': self._create_timeline_slide,
            'diagram': self._create_diagram_slide,
            'progress': self._create_progress_slide,
        }

    def _set_font_style(self, run, style_opts):
        """テキストの一部(run)にフォントスタイルを適用する共通メソッド。"""
        font = run.font
        font.name = style_opts.get('family', CONFIG['FONTS']['family'])
        font.size = Pt(style_opts.get('size', CONFIG['FONTS']['sizes']['body']))
        font.bold = style_opts.get('bold', False)
        color_str = style_opts.get('color', CONFIG['COLORS']['text_primary'])
        font.color.rgb = RGBColor.from_string(color_str)

    def _parse_inline_styles(self, text):
        """`**太字**`や`[[ハイライト]]`のようなインライン記法を解析し、スタイル情報を持つパーツのリストに分解する。"""
        pattern = r'(\*\*|\[\[)(.*?)\1'
        parts, last_end = [], 0
        for match in re.finditer(pattern, text):
            start, end = match.span()
            tag, content = match.groups()
            if start > last_end: parts.append({'text': text[last_end:start], 'style': {}})
            style = {'bold': True, 'color': CONFIG['COLORS']['primary_blue']} if tag == '[[' else {'bold': True}
            parts.append({'text': content, 'style': style})
            last_end = end
        if last_end < len(text): parts.append({'text': text[last_end:], 'style': {}})
        return parts

    def _set_styled_text(self, text_frame, text, base_style_opts):
        """テキストボックスに、インラインスタイルを適用しながらテキストを設定する。"""
        text_frame.clear()
        p = text_frame.paragraphs[0]
        p.font.name = base_style_opts.get('family', CONFIG['FONTS']['family'])
        p.font.size = Pt(base_style_opts.get('size', CONFIG['FONTS']['sizes']['body']))
        p.alignment = base_style_opts.get('align', PP_ALIGN.LEFT)
        for i, line in enumerate(text.split('\n')):
            if i > 0: p = text_frame.add_paragraph()
            parts = self._parse_inline_styles(line) or [{'text': line, 'style': {}}]
            for part in parts:
                run = p.add_run()
                run.text = part['text']
                self._set_font_style(run, {**base_style_opts, **part['style']})

    def _set_bullets_with_inline_styles(self, text_frame, points, base_style_opts):
        """箇条書きリストをテキストボックスに設定する。"""
        text_frame.clear()
        for i, point_text in enumerate(points):
            p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
            run_bullet = p.add_run()
            run_bullet.text = "• "
            self._set_font_style(run_bullet, base_style_opts)
            
            parts = self._parse_inline_styles(point_text)
            for part in parts:
                run = p.add_run()
                run.text = part['text']
                self._set_font_style(run, {**base_style_opts, **part['style']})

    def _draw_bottom_bar(self, slide, layout):
        """スライド下部の青い線を描画する。"""
        rect = layout.get_rect('bottomBar')
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, rect['left'], rect['top'], rect['width'], rect['height'])
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor.from_string(CONFIG['COLORS']['primary_blue'])
        shape.line.fill.background()

    def _add_google_footer(self, slide, layout, page_num):
        """フッター（コピーライトとページ番号）を描画する。"""
        rect = layout.get_rect('footer.leftText')
        textbox = slide.shapes.add_textbox(rect['left'], rect['top'], rect['width'], rect['height'])
        self._set_styled_text(textbox.text_frame, CONFIG['FOOTER_TEXT'], {'size': CONFIG['FONTS']['sizes']['footer']})
        if page_num > 0:
            rect = layout.get_rect('footer.rightPage')
            textbox = slide.shapes.add_textbox(rect['left'], rect['top'], rect['width'], rect['height'])
            self._set_styled_text(textbox.text_frame, str(page_num), {
                'size': CONFIG['FONTS']['sizes']['footer'], 'color': CONFIG['COLORS']['primary_blue'], 'align': PP_ALIGN.RIGHT
            })

    def _draw_standard_title_header(self, slide, layout, key, title):
        """コンテンツスライドの共通ヘッダー（ロゴ、タイトル、下線）を描画する。"""
        rect = layout.get_rect(f'{key}.headerLogo')
        try:
            response = requests.get(CONFIG['LOGOS']['header'])
            slide.shapes.add_picture(io.BytesIO(response.content), rect['left'], rect['top'], width=rect['width'])
        except Exception as e:
            print(f"Warning: ヘッダーロゴの読み込みに失敗しました。 {e}")
        rect = layout.get_rect(f'{key}.title')
        textbox = slide.shapes.add_textbox(rect['left'], rect['top'], rect['width'], rect['height'])
        self._set_styled_text(textbox.text_frame, title, {'size': CONFIG['FONTS']['sizes']['contentTitle'], 'bold': True})
        rect = layout.get_rect(f'{key}.titleUnderline')
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, rect['left'], rect['top'], rect['width'], rect['height'])
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor.from_string(CONFIG['COLORS']['primary_blue'])
        shape.line.fill.background()

    def _draw_subhead_if_any(self, slide, layout, key, subhead):
        """サブヘッド（小見出し）があれば描画する。"""
        if not subhead: return 0
        rect = layout.get_rect(f'{key}.subhead')
        textbox = slide.shapes.add_textbox(rect['left'], rect['top'], rect['width'], rect['height'])
        self._set_styled_text(textbox.text_frame, subhead, {'size': CONFIG['FONTS']['sizes']['subhead']})
        return layout.px_to_emu(36)

    # --- ここから下は、各スライドタイプを描画するための具体的なメソッド群 ---
    # (各メソッドの詳細は省略。内容は元のコードと同じ)

    def _create_title_slide(self, slide, data, layout, page_num):
        """タイプ 'title' のスライド（表紙）を作成します。"""
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor.from_string(CONFIG['COLORS']['background_white'])
        rect = layout.get_rect('titleSlide.logo')
        try:
            response = requests.get(CONFIG['LOGOS']['header'])
            slide.shapes.add_picture(io.BytesIO(response.content), rect['left'], rect['top'], width=rect['width'])
        except Exception as e:
            print(f"Warning: Could not load title logo. {e}")
        rect = layout.get_rect('titleSlide.title')
        textbox = slide.shapes.add_textbox(rect['left'], rect['top'], rect['width'], rect['height'])
        self._set_styled_text(textbox.text_frame, data.get('title', ''), {'size': CONFIG['FONTS']['sizes']['title'], 'bold': True})
        rect = layout.get_rect('titleSlide.date')
        textbox = slide.shapes.add_textbox(rect['left'], rect['top'], rect['width'], rect['height'])
        self._set_styled_text(textbox.text_frame, data.get('date', ''), {'size': CONFIG['FONTS']['sizes']['date']})
        self._draw_bottom_bar(slide, layout)

    def _create_section_slide(self, slide, data, layout, page_num):
        """タイプ 'section' のスライド（章の扉）を作成します。"""
        self._section_counter += 1
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor.from_string(CONFIG['COLORS']['background_gray'])
        num_str = str(data.get('sectionNo', self._section_counter)).zfill(2)
        rect = layout.get_rect('sectionSlide.ghostNum')
        textbox = slide.shapes.add_textbox(rect['left'], rect['top'], rect['width'], rect['height'])
        self._set_styled_text(textbox.text_frame, num_str, {
            'size': CONFIG['FONTS']['sizes']['ghostNum'], 'bold': True, 'color': CONFIG['COLORS']['ghost_gray'], 'align': PP_ALIGN.CENTER
        })
        rect = layout.get_rect('sectionSlide.title')
        textbox = slide.shapes.add_textbox(rect['left'], rect['top'], rect['width'], rect['height'])
        self._set_styled_text(textbox.text_frame, data.get('title', ''), {
            'size': CONFIG['FONTS']['sizes']['sectionTitle'], 'bold': True, 'align': PP_ALIGN.CENTER
        })
        self._add_google_footer(slide, layout, page_num)

    def _create_content_slide(self, slide, data, layout, page_num):
        """タイプ 'content' のスライド（基本的な箇条書き）を作成します。"""
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor.from_string(CONFIG['COLORS']['background_white'])
        self._draw_standard_title_header(slide, layout, 'contentSlide', data.get('title'))
        dy = self._draw_subhead_if_any(slide, layout, 'contentSlide', data.get('subhead'))
        body_rect = layout.get_rect('contentSlide.body')
        body_rect['top'] += dy
        textbox = slide.shapes.add_textbox(body_rect['left'], body_rect['top'], body_rect['width'], body_rect['height'])
        self._set_bullets_with_inline_styles(textbox.text_frame, data.get('points', []), {})
        self._draw_bottom_bar(slide, layout)
        self._add_google_footer(slide, layout, page_num)

    def _create_cards_slide(self, slide, data, layout, page_num):
        """タイプ 'cards' のスライド（カード形式）を作成します。"""
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor.from_string(CONFIG['COLORS']['background_white'])
        self._draw_standard_title_header(slide, layout, 'contentSlide', data.get('title'))
        dy = self._draw_subhead_if_any(slide, layout, 'contentSlide', data.get('subhead'))
        area = layout.get_rect('contentSlide.gridArea')
        area['top'] += dy
        items = data.get('items', [])
        cols = data.get('columns', 3)
        rows = -(-len(items) // cols)
        gap = layout.px_to_emu(16)
        card_w = (area['width'] - gap * (cols - 1)) // cols
        card_h = (area['height'] - gap * (rows - 1)) // rows
        for i, item in enumerate(items):
            r, c = divmod(i, cols)
            left, top = area['left'] + c * (card_w + gap), area['top'] + r * (card_h + gap)
            shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, card_w, card_h)
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor.from_string(CONFIG['COLORS']['card_bg'])
            shape.line.color.rgb = RGBColor.from_string(CONFIG['COLORS']['card_border'])
            shape.line.width = Pt(1).emu
            tf = shape.text_frame
            tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = layout.px_to_emu(10)
            tf.word_wrap = True
            tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
            if isinstance(item, dict):
                self._set_styled_text(tf, f"**{item.get('title', '')}**\n{item.get('desc', '')}", {'size': CONFIG['FONTS']['sizes']['cardDesc']})
        self._draw_bottom_bar(slide, layout)
        self._add_google_footer(slide, layout, page_num)

    def _create_table_slide(self, slide, data, layout, page_num):
        """タイプ 'table' のスライド（表）を作成します。"""
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor.from_string(CONFIG['COLORS']['background_white'])
        self._draw_standard_title_header(slide, layout, 'contentSlide', data.get('title'))
        dy = self._draw_subhead_if_any(slide, layout, 'contentSlide', data.get('subhead'))
        area = layout.get_rect('contentSlide.body')
        area['top'] += dy
        headers = data.get('headers', [])
        rows_data = data.get('rows', [])
        if not headers or not rows_data: return
        shape = slide.shapes.add_table(len(rows_data) + 1, len(headers), area['left'], area['top'], area['width'], area['height'])
        table = shape.table
        for c, header_text in enumerate(headers):
            cell = table.cell(0, c)
            cell.text = header_text
            cell.text_frame.paragraphs[0].font.bold = True
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        for r, row_data in enumerate(rows_data):
            for c, cell_text in enumerate(row_data):
                cell = table.cell(r + 1, c)
                cell.text = cell_text
                cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        self._draw_bottom_bar(slide, layout)
        self._add_google_footer(slide, layout, page_num)

    def _create_compare_slide(self, slide, data, layout, page_num):
        """タイプ 'compare' のスライド（2項目比較）を作成します。"""
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor.from_string(CONFIG['COLORS']['background_white'])
        self._draw_standard_title_header(slide, layout, 'contentSlide', data.get('title'))
        dy = self._draw_subhead_if_any(slide, layout, 'contentSlide', data.get('subhead'))
        def draw_compare_box(rect, title, items):
            slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, rect['left'], rect['top'], rect['width'], rect['height'])
            title_h = layout.px_to_emu(40)
            title_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, rect['left'], rect['top'], rect['width'], title_h)
            title_bar.fill.solid()
            title_bar.fill.fore_color.rgb = RGBColor.from_string(CONFIG['COLORS']['primary_blue'])
            title_bar.line.fill.background()
            title_bar.text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
            self._set_styled_text(title_bar.text_frame, title, {'size': CONFIG['FONTS']['sizes']['laneTitle'], 'bold': True, 'color': CONFIG['COLORS']['text_white'], 'align': PP_ALIGN.CENTER})
            pad = layout.px_to_emu(12)
            text_box = slide.shapes.add_textbox(rect['left'] + pad, rect['top'] + title_h + pad, rect['width'] - 2*pad, rect['height'] - title_h - 2*pad)
            self._set_bullets_with_inline_styles(text_box.text_frame, items, {})
        left_rect = layout.get_rect('contentSlide.compareLeft')
        left_rect['top'] += dy
        draw_compare_box(left_rect, data.get('leftTitle', ''), data.get('leftItems', []))
        right_rect = layout.get_rect('contentSlide.compareRight')
        right_rect['top'] += dy
        draw_compare_box(right_rect, data.get('rightTitle', ''), data.get('rightItems', []))
        self._draw_bottom_bar(slide, layout)
        self._add_google_footer(slide, layout, page_num)

    def _create_process_slide(self, slide, data, layout, page_num):
        """タイプ 'process' のスライド（手順・工程）を作成します。"""
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor.from_string(CONFIG['COLORS']['background_white'])
        self._draw_standard_title_header(slide, layout, 'contentSlide', data.get('title'))
        dy = self._draw_subhead_if_any(slide, layout, 'contentSlide', data.get('subhead'))
        area = layout.get_rect('contentSlide.body')
        area['top'] += dy
        steps = data.get('steps', [])
        n = len(steps)
        if n == 0: return
        gap_y = (area['height'] - layout.px_to_emu(40)) / max(1, n - 1) if n > 1 else 0
        cx = area['left'] + layout.px_to_emu(44)
        top0 = area['top'] + layout.px_to_emu(20)
        if n > 1:
            line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, cx - layout.px_to_emu(1), top0, layout.px_to_emu(2), gap_y * (n - 1))
            line.fill.solid()
            line.fill.fore_color.rgb = RGBColor.from_string(CONFIG['COLORS']['faint_gray'])
            line.line.fill.background()
        for i, step_text in enumerate(steps):
            cy = top0 + gap_y * i
            sz = layout.px_to_emu(28)
            num_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, cx - sz/2, cy - sz/2, sz, sz)
            num_box.fill.solid()
            num_box.fill.fore_color.rgb = RGBColor.from_string(CONFIG['COLORS']['background_white'])
            num_box.line.color.rgb = RGBColor.from_string(CONFIG['COLORS']['primary_blue'])
            num_box.line.width = Pt(1).emu
            num_box.text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
            self._set_styled_text(num_box.text_frame, str(i + 1), {'size': 12, 'bold': True, 'color': CONFIG['COLORS']['primary_blue'], 'align': PP_ALIGN.CENTER})
            txt_box = slide.shapes.add_textbox(cx + layout.px_to_emu(28), cy - layout.px_to_emu(16), area['width'] - layout.px_to_emu(100), layout.px_to_emu(32))
            txt_box.text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
            self._set_styled_text(txt_box.text_frame, step_text, {'size': CONFIG['FONTS']['sizes']['processStep']})
        self._draw_bottom_bar(slide, layout)
        self._add_google_footer(slide, layout, page_num)

    def _create_timeline_slide(self, slide, data, layout, page_num):
        """タイプ 'timeline' のスライド（時系列）を作成します。"""
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor.from_string(CONFIG['COLORS']['background_white'])
        self._draw_standard_title_header(slide, layout, 'contentSlide', data.get('title'))
        dy = self._draw_subhead_if_any(slide, layout, 'contentSlide', data.get('subhead'))
        area = layout.get_rect('contentSlide.body')
        area['top'] += dy
        milestones = data.get('milestones', [])
        n = len(milestones)
        if n == 0: return
        inner_margin = layout.px_to_emu(60)
        base_y = area['top'] + area['height'] * 0.5
        left_x, right_x = area['left'] + inner_margin, area['left'] + area['width'] - inner_margin
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left_x, base_y - layout.px_to_emu(1), right_x - left_x, layout.px_to_emu(2))
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor.from_string(CONFIG['COLORS']['faint_gray'])
        line.line.fill.background()
        dot_r = layout.px_to_emu(12)
        gap_x = (right_x - left_x) / max(1, n - 1) if n > 1 else 0
        for i, m in enumerate(milestones):
            x = left_x + gap_x * i
            dot = slide.shapes.add_shape(MSO_SHAPE.ELLIPSE, x - dot_r/2, base_y - dot_r/2, dot_r, dot_r)
            state = m.get('state', 'todo').lower()
            if state == 'done':
                dot.fill.solid(); dot.fill.fore_color.rgb = RGBColor.from_string(CONFIG['COLORS']['google_green']); dot.line.fill.background()
            elif state == 'next':
                dot.fill.solid(); dot.fill.fore_color.rgb = RGBColor.from_string(CONFIG['COLORS']['background_white']); dot.line.color.rgb = RGBColor.from_string(CONFIG['COLORS']['google_yellow']); dot.line.width = Pt(2).emu
            else:
                dot.fill.solid(); dot.fill.fore_color.rgb = RGBColor.from_string(CONFIG['COLORS']['background_white']); dot.line.color.rgb = RGBColor.from_string(CONFIG['COLORS']['neutral_gray']); dot.line.width = Pt(1).emu
            label_box = slide.shapes.add_textbox(x - layout.px_to_emu(50), base_y - layout.px_to_emu(50), layout.px_to_emu(100), layout.px_to_emu(30))
            self._set_styled_text(label_box.text_frame, m.get('label', ''), {'size': CONFIG['FONTS']['sizes']['timelineLabel'], 'bold': True, 'align': PP_ALIGN.CENTER})
            date_box = slide.shapes.add_textbox(x - layout.px_to_emu(50), base_y + layout.px_to_emu(15), layout.px_to_emu(100), layout.px_to_emu(20))
            self._set_styled_text(date_box.text_frame, m.get('date', ''), {'size': CONFIG['FONTS']['sizes']['timelineLabel'], 'color': CONFIG['COLORS']['neutral_gray'], 'align': PP_ALIGN.CENTER})
        self._draw_bottom_bar(slide, layout)
        self._add_google_footer(slide, layout, page_num)

    def _create_diagram_slide(self, slide, data, layout, page_num):
        """タイプ 'diagram' のスライド（レーン図）を作成します。"""
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor.from_string(CONFIG['COLORS']['background_white'])
        self._draw_standard_title_header(slide, layout, 'contentSlide', data.get('title'))
        dy = self._draw_subhead_if_any(slide, layout, 'contentSlide', data.get('subhead'))
        area = layout.get_rect('contentSlide.body')
        area['top'] += dy
        lanes = data.get('lanes', [])
        n = len(lanes)
        if n == 0: return
        lane_gap, lane_pad, lane_title_h = layout.px_to_emu(24), layout.px_to_emu(10), layout.px_to_emu(30)
        card_gap, card_min_h = layout.px_to_emu(12), layout.px_to_emu(48)
        arrow_h, arrow_gap = layout.px_to_emu(10), layout.px_to_emu(8)
        lane_w = (area['width'] - lane_gap * (n - 1)) / n
        card_boxes = [[] for _ in range(n)]
        for j, lane in enumerate(lanes):
            left = area['left'] + j * (lane_w + lane_gap)
            lt = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, area['top'], lane_w, lane_title_h)
            lt.fill.solid(); lt.fill.fore_color.rgb = RGBColor.from_string(CONFIG['COLORS']['lane_title_bg']); lt.line.color.rgb = RGBColor.from_string(CONFIG['COLORS']['lane_border'])
            lt.text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
            self._set_styled_text(lt.text_frame, lane.get('title', ''), {'size': CONFIG['FONTS']['sizes']['laneTitle'], 'bold': True, 'align': PP_ALIGN.CENTER})
            items = lane.get('items', [])
            rows = len(items)
            avail_h = area['height'] - lane_title_h - lane_pad * 2
            card_h = max(card_min_h, (avail_h - card_gap * (rows - 1)) / rows if rows > 0 else 0)
            for i, item_text in enumerate(items):
                card_top = area['top'] + lane_title_h + lane_pad + i * (card_h + card_gap)
                card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left + lane_pad, card_top, lane_w - 2 * lane_pad, card_h)
                card.fill.solid(); card.fill.fore_color.rgb = RGBColor.from_string(CONFIG['COLORS']['card_bg']); card.line.color.rgb = RGBColor.from_string(CONFIG['COLORS']['card_border'])
                card.text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
                card.text_frame.margin_left = card.text_frame.margin_right = layout.px_to_emu(8)
                card.text_frame.word_wrap = True
                self._set_styled_text(card.text_frame, item_text, {'size': CONFIG['FONTS']['sizes']['body'], 'align': PP_ALIGN.CENTER})
                card_boxes[j].append({'left': left + lane_pad, 'top': card_top, 'width': lane_w - 2 * lane_pad, 'height': card_h})
        max_rows = max(len(cb) for cb in card_boxes) if card_boxes else 0
        for i in range(max_rows):
            for j in range(n - 1):
                if i < len(card_boxes[j]) and i < len(card_boxes[j+1]):
                    rect_a, rect_b = card_boxes[j][i], card_boxes[j+1][i]
                    from_x, to_x = rect_a['left'] + rect_a['width'], rect_b['left']
                    y_mid = rect_a['top'] + rect_a['height'] / 2
                    arrow_left, arrow_width = from_x + arrow_gap, to_x - from_x - 2 * arrow_gap
                    if arrow_width > 0:
                        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, arrow_left, y_mid - arrow_h/2, arrow_width, arrow_h)
                        arrow.fill.solid(); arrow.fill.fore_color.rgb = RGBColor.from_string(CONFIG['COLORS']['primary_blue']); arrow.line.fill.background()
        self._draw_bottom_bar(slide, layout)
        self._add_google_footer(slide, layout, page_num)

    def _create_progress_slide(self, slide, data, layout, page_num):
        """タイプ 'progress' のスライド（進捗バー）を作成します。"""
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor.from_string(CONFIG['COLORS']['background_white'])
        self._draw_standard_title_header(slide, layout, 'contentSlide', data.get('title'))
        dy = self._draw_subhead_if_any(slide, layout, 'contentSlide', data.get('subhead'))
        area = layout.get_rect('contentSlide.body')
        area['top'] += dy
        items = data.get('items', [])
        n = len(items)
        if n == 0: return
        row_h = area['height'] / n
        for i, item in enumerate(items):
            y = area['top'] + i * row_h + row_h/2 - layout.px_to_emu(9)
            label_box = slide.shapes.add_textbox(area['left'], y, layout.px_to_emu(150), layout.px_to_emu(18))
            self._set_styled_text(label_box.text_frame, item.get('label', ''), {})
            bar_left, bar_w, bar_h = area['left'] + layout.px_to_emu(160), area['width'] - layout.px_to_emu(220), layout.px_to_emu(14)
            bg_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, bar_left, y, bar_w, bar_h)
            bg_bar.fill.solid(); bg_bar.fill.fore_color.rgb = RGBColor.from_string(CONFIG['COLORS']['faint_gray']); bg_bar.line.fill.background()
            percent = max(0, min(100, item.get('percent', 0)))
            fg_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, bar_left, y, bar_w * (percent/100), bar_h)
            fg_bar.fill.solid(); fg_bar.fill.fore_color.rgb = RGBColor.from_string(CONFIG['COLORS']['google_green']); fg_bar.line.fill.background()
            pct_box = slide.shapes.add_textbox(bar_left + bar_w + layout.px_to_emu(6), y - layout.px_to_emu(1), layout.px_to_emu(50), layout.px_to_emu(16))
            self._set_styled_text(pct_box.text_frame, f"{percent}%", {'size': CONFIG['FONTS']['sizes']['timelineLabel'], 'color': CONFIG['COLORS']['neutral_gray']})
        self._draw_bottom_bar(slide, layout)
        self._add_google_footer(slide, layout, page_num)

    def _create_closing_slide(self, slide, data, layout, page_num):
        """タイプ 'closing' のスライド（結び）を作成します。"""
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor.from_string(CONFIG['COLORS']['background_white'])
        try:
            response = requests.get(CONFIG['LOGOS']['closing'])
            logo_data = io.BytesIO(response.content)
            with Image.open(logo_data) as img:
                aspect_ratio = img.height / img.width
            logo_w_emu = layout.px_to_emu(450)
            logo_h_emu = int(logo_w_emu * aspect_ratio)
            left, top = (layout.page_w_emu - logo_w_emu) / 2, (layout.page_h_emu - logo_h_emu) / 2
            logo_data.seek(0)
            slide.shapes.add_picture(logo_data, left, top, width=logo_w_emu)
        except Exception as e:
            print(f"Warning: Could not load closing logo. {e}")

# ==============================================================================
# 3. PowerPoint生成のメイン関数
# ==============================================================================
def generate_presentation_in_memory(data):
    """
    slide_dataを元にPowerPointプレゼンテーションをメモリ上に生成し、
    そのバイトデータを返す。これにより、Lambda環境でファイルシステムに書き込まずに済む。
    """
    # 新しいプレゼンテーションオブジェクトを作成
    prs = Presentation()
    # スライドサイズを設定
    prs.slide_width, prs.slide_height = CONFIG['BASE_EMU']['W'], CONFIG['BASE_EMU']['H']
    
    # レイアウトマネージャーとスライドジェネレーターを初期化
    layout_manager = LayoutManager(prs.slide_width, prs.slide_height)
    generator = SlideGenerator()
    
    page_counter = 0
    # slide_dataの各要素をループ処理
    for item in data:
        # 表紙と結びのスライド以外でページ番号をカウントアップ
        if item['type'] not in ['title', 'closing']: page_counter += 1
        
        # 対応するスライド生成関数を取得
        generator_func = generator.slide_generators.get(item['type'])
        if generator_func:
            # 白紙のスライドを追加
            slide = prs.slides.add_slide(prs.slide_layouts[6]) 
            # スライド生成関数を実行
            generator_func(slide, item, layout_manager, page_counter)
            # スピーカーノートがあれば追加
            if item.get('notes'):
                slide.notes_slide.notes_text_frame.text = item['notes']
    
    # ファイルに保存する代わりに、メモリ上のバイトストリームに保存
    powerpoint_stream = io.BytesIO()
    prs.save(powerpoint_stream)
    powerpoint_stream.seek(0) # ストリームの先頭に戻す
    
    return powerpoint_stream.getvalue()

# ==============================================================================
# 4. AWS Lambda ハンドラ関数
# ==============================================================================
def lambda_handler(event, context):
    """
    AWS Lambdaのエントリーポイント。API GatewayからのHTTPリクエストを処理する。
    `event`引数にはリクエスト情報（ヘッダー、ボディなど）が、
    `context`引数には実行環境情報が含まれる。
    """
    print("--- Lambda関数がトリガーされました ---")
    
    try:
        # 環境変数からS3バケット名を取得。これはLambdaのコンソールで設定する必要がある。
        s3_bucket_name = os.environ.get('S3_BUCKET_NAME')
        if not s3_bucket_name:
            # バケット名が設定されていない場合は、設定不備としてエラーを発生させる
            raise ValueError("環境変数 'S3_BUCKET_NAME' が設定されていません。")

        # API Gatewayからのリクエストボディを取得・パース
        print("リクエストボディを解析します...")
        if 'body' not in event or not event['body']:
             raise ValueError("リクエストボディが空です。")
        
        body = json.loads(event['body'])
        slide_data_string = body.get('slideData')
        
        if not slide_data_string:
            raise ValueError("リクエストボディに 'slideData' が含まれていません。")
        print("'slideData'の取得に成功しました。")

        # slideData文字列をPythonのリストオブジェクトに安全に変換
        try:
            # ast.literal_evalは、eval()よりも安全に文字列をPythonのデータ構造に変換する
            slide_data_list = ast.literal_eval(slide_data_string.strip())
            if not isinstance(slide_data_list, list):
                raise TypeError()
        except (ValueError, SyntaxError, TypeError):
            raise ValueError("slideDataの形式が不正です。Pythonのリスト形式の文字列である必要があります。")
        print("slideDataをPythonリストに変換しました。")

        # PowerPointファイルをメモリ上にバイトデータとして生成
        print("PowerPointファイルの生成を開始します...")
        powerpoint_bytes = generate_presentation_in_memory(slide_data_list)
        print("PowerPointファイルの生成が完了しました。")

        # 一意なファイル名を生成してS3にアップロード
        # UUIDを使うことで、ファイル名の衝突を避ける
        file_name = f"presentations/{uuid.uuid4()}.pptx"
        print(f"S3バケット '{s3_bucket_name}' に '{file_name}' としてアップロードします...")
        s3_client.put_object(
            Bucket=s3_bucket_name,
            Key=file_name,
            Body=powerpoint_bytes,
            ContentType='application/vnd.openxmlformats-officedocument.presentationml.presentation'
        )
        print("S3へのアップロードが完了しました。")

        # S3オブジェクトの署名付きURLを生成（有効期限: 1時間）
        # これにより、プライベートなS3オブジェクトに一時的な公開アクセスを許可できる
        print("署名付きダウンロードURLを生成します...")
        try:
            presigned_url = s3_client.generate_presigned_url(
                'get_object',
                Params={'Bucket': s3_bucket_name, 'Key': file_name},
                ExpiresIn=3600  # 3600秒 = 1時間
            )
            print("URLの生成に成功しました。")
        except ClientError as e:
            print(f"署名付きURLの生成に失敗しました: {e}")
            raise

        # 成功レスポンスを返す
        # このレスポンスがAPI Gatewayを通じて呼び出し元（Salesforce）に返される
        return {
            'statusCode': 200,
            'headers': {
                'Content-Type': 'application/json',
                'Access-Control-Allow-Origin': '*' # CORSを許可（本番環境ではドメインを限定することが望ましい）
            },
            'body': json.dumps({
                'message': 'プレゼンテーションが正常に作成されました。',
                'downloadUrl': presigned_url,
                's3Key': file_name  # Salesforceで保存・再利用するためのS3キーも返す
            })
        }

    except Exception as e:
        # エラーハンドリング: 処理中に発生した例外をキャッチする
        print(f"エラーが発生しました: {e}")
        # エラーの種類に応じてステータスコードを分ける
        error_type = "BadRequest" if isinstance(e, (ValueError, KeyError, TypeError)) else "InternalServerError"
        status_code = 400 if error_type == "BadRequest" else 500
        
        return {
            'statusCode': status_code,
            'headers': {
                'Content-Type': 'application/json',
                'Access-Control-Allow-Origin': '*'
            },
            'body': json.dumps({'error': str(e)})
        }
